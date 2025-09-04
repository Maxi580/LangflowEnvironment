from langflow.custom import Component
from langflow.io import Output, MultilineInput, SecretStrInput
from langflow.schema import Message
from pptx.util import Inches
import google.generativeai as genai
from pptx.dml.color import RGBColor
from PIL import Image
import io
import base64
import tempfile
import os
from pptx import Presentation
from typing import List, Dict, Tuple, Any
from dotenv import load_dotenv

load_dotenv()

PPTX_MAGIC_BYTES = os.getenv("PPTX_MAGIC_BYTES")
TEMPLATE_PATH = os.getenv("TEMPLATE_PATH")


class CombinedPPTXExtractorCreator(Component):
    display_name = "Case Study Library Transformator"
    description = "Extracts Challenge/Solution/Value from PPTX and directly creates PowerPoint presentations"
    icon = "🔄"

    inputs = [
        MultilineInput(
            name="b64_data",
            display_name="Base64 PPTX Data",
            info="Comma-separated base64 encoded PPTX file data",
            required=True,
        ),
        SecretStrInput(name="api_key"),
    ]

    outputs = [
        Output(
            display_name="PowerPoint Files",
            name="powerpoint_output",
            method="extract_and_create_powerpoints"
        )
    ]

    def is_pptx_file(self, file_data: bytes) -> bool:
        """Check if the file data represents a valid PPTX file"""
        try:
            if not file_data.startswith(b'PK\x03\x04'):
                return False

            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_file.write(file_data)
                temp_file_path = temp_file.name

            try:
                prs = Presentation(temp_file_path)
                _ = len(prs.slides)
                return True
            except Exception:
                return False
            finally:
                try:
                    os.unlink(temp_file_path)
                except:
                    pass

        except Exception:
            return False

    def get_text_shapes_from_slide(self, slide) -> List[Dict]:
        """Extract all text shapes from a slide with their spatial information"""
        text_shapes = []

        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text = shape.text_frame.text.strip()
                text_shapes.append({
                    'text': text,
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })

        return text_shapes

    def find_text_below_title(self, text_shapes: List[Dict], title_keywords: List[str], x_margin: int = 720000) -> str:
        """
        Find a title with any of the keywords (in order) and return the first text field below it with similar x-coordinate
        """
        # Try each keyword in order until we find a match
        for keyword in title_keywords:
            title_shape = None

            # Find the title shape containing this specific keyword
            for shape in text_shapes:
                text_lower = shape['text'].lower().strip()
                if keyword.lower() in text_lower and len(shape['text']) <= 50:  # Likely a title
                    title_shape = shape
                    break

            # If we found a title with this keyword, look for content below it
            if title_shape:
                title_x = title_shape['left']
                title_y = title_shape['top']

                # Find text fields below the title with similar x-coordinate
                candidates = []

                for shape in text_shapes:
                    # Skip the title itself
                    if shape['left'] == title_x and shape['top'] == title_y:
                        continue

                    # Must be below the title
                    if shape['top'] <= title_y:
                        continue

                    # Check if x-coordinate is within margin
                    x_distance = abs(shape['left'] - title_x)
                    if x_distance <= x_margin:
                        y_distance = shape['top'] - title_y
                        candidates.append({
                            'text': shape['text'],
                            'y_distance': y_distance,
                            'x_distance': x_distance
                        })

                if candidates:
                    # Sort by y_distance (closest below the title) then by x_distance
                    candidates.sort(key=lambda x: (x['y_distance'], x['x_distance']))
                    return candidates[0]['text']

        return ""

    def find_project_name(self, text_shapes: List[Dict], target_x_cm: float = 1.19, margin_cm: float = 0.2) -> str:
        """
        Find the project name at horizontal position 1.19cm - it's the second highest text element in that area
        """
        # Convert cm to EMU (English Metric Units used by python-pptx)
        # 1 cm = 360000 EMU
        target_x_emu = int(target_x_cm * 360000)
        margin_emu = int(margin_cm * 360000)

        # Find all text shapes at the target x position (within margin)
        candidates = []

        for shape in text_shapes:
            x_distance = abs(shape['left'] - target_x_emu)
            if x_distance <= margin_emu:
                candidates.append({
                    'text': shape['text'],
                    'top': shape['top'],
                    'y_position': shape['top']
                })

        if len(candidates) < 2:
            return ""  # Need at least 2 elements (sector and project name)

        # Sort by y position (top to bottom)
        candidates.sort(key=lambda x: x['y_position'])

        # Return the second highest (index 1) - this should be the project name
        return candidates[1]['text']

    def find_logo_in_area(self, slide, target_x_cm: float = 21.98, target_y_cm: float = 0.46,
                          width_cm: float = 11.88, height_cm: float = 8.1) -> str:
        """
        Find and extract logo from the specified area and return as base64 encoded string
        """
        # Convert cm to EMU (English Metric Units used by python-pptx)
        # 1 cm = 360000 EMU
        target_x_emu = int(target_x_cm * 360000)
        target_y_emu = int(target_y_cm * 360000)
        width_emu = int(width_cm * 360000)
        height_emu = int(height_cm * 360000)

        # Define the search area bounds
        left_bound = target_x_emu
        right_bound = target_x_emu + width_emu
        top_bound = target_y_emu
        bottom_bound = target_y_emu + height_emu

        for shape in slide.shapes:
            # Check if shape is an image/picture
            if hasattr(shape, 'image'):
                # Check if the shape is within our target area
                shape_left = shape.left
                shape_top = shape.top
                shape_right = shape.left + shape.width
                shape_bottom = shape.top + shape.height

                # Check if shape overlaps with our target area
                if (shape_left < right_bound and shape_right > left_bound and
                        shape_top < bottom_bound and shape_bottom > top_bound):

                    try:
                        # Extract the image data
                        image = shape.image
                        image_bytes = image.blob

                        # Convert to base64
                        import base64
                        base64_content = base64.b64encode(image_bytes).decode('utf-8')
                        return base64_content

                    except Exception as e:
                        # If extraction fails, continue to next shape
                        continue

        return ""  # No logo found in the specified area

    def analyze_client_agent(self, logo_base64: str, challenge: str, solution: str, business_impact: str,
                             project_name: str) -> dict:
        if not self.api_key:
            return {"customer_name": "Unknown Client", "about_client": "API key not provided"}

        try:
            genai.configure(api_key=self.api_key)
            model = genai.GenerativeModel("gemini-2.5-flash")

            content = []

            if logo_base64:
                logo_bytes = base64.b64decode(logo_base64)
                content.append({
                    "mime_type": "image/png",
                    "data": logo_bytes
                })

            prompt = f"""
            Your Job is to research Companies. You want to create small Informational texts about them that contain the following Information:

            Based on this company logo (if provided) and project details:
            Project: {project_name}
            Challenge: {challenge}
            Solution: {solution}
            Business Impact: {business_impact}

            Research and provide:
            1. COMPANY_NAME: Full company name (research full name if only acronym visible)
            2. ABOUT_CLIENT: Concise description (75-100 tokens) covering:
               - Basic Information: Type of organization, subsidiaries/ownership
               - Business Model: Core business, market position, geographic operations
               - Financials: Annual revenue (€), employee count
               - Transformation: Industry challenges/trends

            Format your response exactly as:
            COMPANY_NAME: [company name]
            ABOUT_CLIENT: [75-100 token description in continuous text, no line breaks]
            """

            content.append(prompt)

            response = model.generate_content(content)
            return self.parse_ai_response(response.text)

        except Exception as e:
            return {"customer_name": "Unknown Client", "about_client": f"Analysis failed: {str(e)}"}

    def parse_ai_response(self, response_text: str) -> dict:
        try:
            lines = response_text.strip().split('\n')
            company_name = "Unknown Client"
            about_client = "Company analysis not available"

            for line in lines:
                if line.startswith("COMPANY_NAME:"):
                    company_name = line.replace("COMPANY_NAME:", "").strip()
                elif line.startswith("ABOUT_CLIENT:"):
                    about_client = line.replace("ABOUT_CLIENT:", "").strip()

            return {
                "customer_name": company_name,
                "about_client": about_client
            }
        except:
            return {"customer_name": "Unknown Client", "about_client": "Parsing failed"}

    def extract_fields_from_slide(self, slide, slide_number: int) -> Dict[str, str]:
        """Extract Challenge, Solution, and Value from a single slide"""
        text_shapes = self.get_text_shapes_from_slide(slide)

        challenge = self.find_text_below_title(text_shapes, ["Challenge"])
        solution = self.find_text_below_title(text_shapes, ["Solution"])
        business_impact = self.find_text_below_title(text_shapes, ["Value", "Business Benefits"])
        project_name = self.find_project_name(text_shapes)
        logo_base64 = self.find_logo_in_area(slide)
        analysis_result = self.analyze_client_agent(logo_base64, challenge, solution, business_impact, project_name)

        return {
            'slide_number': slide_number,
            'customer_name': analysis_result.get('customer_name', f"Unknown Client at Reference {slide_number}"),
            'project_name': project_name,
            'about_client': analysis_result.get('about_client', "Client information extracted from presentation"),
            'challenge_text': challenge,
            'solution_text': solution,
            'impact_text': business_impact,
            'why_us_text': '',
            'it_impact_text': '',
            'methods_text': '',
            'software_used_impact': '',
            'hardware_text': '',
            'network_communication_text': '',
            'technology_used_impact': '',
            'logo_base64': logo_base64
        }

    def has_valid_content(self, slide_data: Dict[str, str]) -> str:
        """Check if slide has any extractable content"""
        return slide_data['challenge_text'] or slide_data['solution_text'] or slide_data['impact_text']

    def find_and_replace_text_in_slide(self, slide, replacements):
        """Find and replace text in all text boxes on a slide"""
        replacements_made = 0

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, replacement in replacements.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, replacement)
                                replacements_made += 1

            elif hasattr(shape, "table"):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for placeholder, replacement in replacements.items():
                                    if placeholder in run.text:
                                        run.text = run.text.replace(placeholder, replacement)
                                        replacements_made += 1

        return replacements_made

    def decode_base64_image(self, base64_string):
        """
        Decode base64 image string and return image bytes
        """
        try:
            if base64_string.startswith('data:'):
                base64_string = base64_string.split(',', 1)[1]

            image_data = base64.b64decode(base64_string)

            img = Image.open(io.BytesIO(image_data))

            # Convert to RGB if necessary (for JPEG compatibility)
            if img.mode in ('RGBA', 'LA', 'P'):
                # Create white background
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                img = background

            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            return img_buffer.getvalue(), img.size

        except Exception as e:
            print(f"Error decoding base64 image: {e}")
            return None, None

    def add_logo_to_slide(self, slide, logo_base64):
        """
        Add logo to slide at fixed position with aspect ratio preservation
        """
        if not logo_base64:
            return False

        try:
            logo_data, logo_size = self.decode_base64_image(logo_base64)
            if not logo_data:
                return False

            # Convert cm to inches (PowerPoint uses inches)
            cm_to_inches = 0.393701

            # Available space dimensions
            max_width = Inches(2.87 * cm_to_inches)
            max_height = Inches(2.53 * cm_to_inches)

            # Position
            left = Inches(29.81 * cm_to_inches)
            top = Inches(0.81 * cm_to_inches)

            # Calculate aspect ratio preserving dimensions
            original_width, original_height = logo_size
            aspect_ratio = original_width / original_height

            # Fit within the available space while maintaining aspect ratio
            if aspect_ratio > (max_width / max_height):
                # Image is wider - fit to width
                final_width = max_width
                final_height = max_width / aspect_ratio
            else:
                # Image is taller - fit to height
                final_height = max_height
                final_width = max_height * aspect_ratio

            # Center the image within the available space
            centered_left = left + (max_width - final_width) / 2
            centered_top = top + (max_height - final_height) / 2

            # Save logo to temporary file
            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_file.write(logo_data)
                tmp_logo_path = tmp_file.name

            try:
                # Add the logo to the slide with preserved aspect ratio
                pic = slide.shapes.add_picture(tmp_logo_path, centered_left, centered_top, final_width, final_height)
                return True
            finally:
                try:
                    os.unlink(tmp_logo_path)
                except:
                    pass

        except Exception as e:
            print(f"Error adding logo to slide: {e}")
            return False

    def create_powerpoint_from_data(self, reference_data: Dict[str, str], reference_index: int) -> str:
        """Create a single PowerPoint from extracted reference data using simple template"""
        try:
            if not os.path.exists(TEMPLATE_PATH):
                return f"Template file not found at {TEMPLATE_PATH}"

            prs = Presentation(TEMPLATE_PATH)

            replacements = {
                '{{CUSTOMER_NAME}}': reference_data['customer_name'],
                '{{ABOUT_CLIENT}}': reference_data['about_client'],
                '{{PROJECT_NAME}}': reference_data['project_name'],
                '{{CHALLENGE_TEXT}}': reference_data['challenge_text'],
                '{{SOLUTION_TEXT}}': reference_data['solution_text'],
                '{{IMPACT_TEXT}}': reference_data['impact_text'],
            }

            for slide_idx, slide in enumerate(prs.slides):
                self.find_and_replace_text_in_slide(slide, replacements)

                # Add logo to first slide if available
                if slide_idx == 0 and reference_data['logo_base64']:
                    self.add_logo_to_slide(slide, reference_data['logo_base64'])

            # Save to temporary file
            temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            temp_file_path = temp_file.name
            temp_file.close()

            prs.save(temp_file_path)

            # Read and encode
            with open(temp_file_path, 'rb') as f:
                file_content = f.read()

            os.unlink(temp_file_path)

            # Encode to base64
            base64_content = base64.b64encode(file_content).decode('utf-8')

            client_name = reference_data['customer_name']
            safe_client_name = "".join(c for c in client_name if c.isalnum() or c in (' ', '-', '_')).rstrip()
            safe_client_name = safe_client_name.replace(' ', '_')
            if safe_client_name and safe_client_name != "Unknown_Client":
                filename = f"{safe_client_name}_slide_{reference_data['slide_number']}.pptx"
            else:
                filename = f"reference_{reference_index}_slide_{reference_data['slide_number']}.pptx"

            return f"""Reference {reference_index} PowerPoint created successfully!
    <{PPTX_MAGIC_BYTES}>
    filename:{filename}
    content_type:application/vnd.openxmlformats-officedocument.presentationml.presentation
    size:{len(file_content)}
    data:{base64_content}
    </{PPTX_MAGIC_BYTES}>

    """

        except Exception as e:
            return f"Failed to create PowerPoint for reference {reference_index}: {str(e)}\n"

    def extract_and_create_from_pptx(self, file_data: bytes, file_index: int) -> Tuple[bool, str, int]:
        """Extract data from PPTX and create PowerPoint files for each valid reference"""
        try:
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_file.write(file_data)
                temp_file_path = temp_file.name

            try:
                prs = Presentation(temp_file_path)
                output_text = f"📊 Processing File {file_index}:\n"
                powerpoints_created = 0
                reference_index = 1

                # Process each slide
                for slide_idx, slide in enumerate(prs.slides, 1):
                    slide_data = self.extract_fields_from_slide(slide, slide_idx)

                    # Only create PowerPoint if slide has valid content
                    if self.has_valid_content(slide_data):
                        powerpoint_result = self.create_powerpoint_from_data(slide_data, reference_index)
                        output_text += powerpoint_result
                        powerpoints_created += 1
                        reference_index += 1

                output_text += f"✅ File {file_index} processed: {powerpoints_created} PowerPoints created from {len(prs.slides)} slides\n\n"
                return True, output_text, powerpoints_created

            finally:
                try:
                    os.unlink(temp_file_path)
                except:
                    pass

        except Exception as e:
            return False, f"❌ Error processing file {file_index}: {str(e)}\n", 0

    def clean_b64_string(self, b64_string: str) -> str:
        """Clean base64 string by removing data URI prefixes and whitespace"""
        b64_string = b64_string.strip()

        if b64_string.startswith('data:'):
            if ',' in b64_string:
                b64_string = b64_string.split(',', 1)[1]

        b64_string = ''.join(b64_string.split())
        return b64_string

    def process_single_pptx(self, b64_string: str, index: int) -> Tuple[str, bool, str, int]:
        """Process a single PPTX file from base64 string"""
        try:
            clean_b64 = self.clean_b64_string(b64_string)

            if not clean_b64:
                return f"⚠️ Empty base64 string for file {index}", False, "", 0

            file_data = base64.b64decode(clean_b64)

            if not self.is_pptx_file(file_data):
                return f"❌ File {index} is not a valid PPTX format", False, "", 0

            success, content, count = self.extract_and_create_from_pptx(file_data, index)

            if success:
                return f"✅ Successfully processed PPTX file {index}", True, content, count
            else:
                return f"❌ Error processing file {index}: {content}", False, "", 0

        except Exception as e:
            error_msg = f"❌ Error processing file {index}: {str(e)}"
            return error_msg, False, "", 0

    def extract_and_create_powerpoints(self) -> Message:
        """Main method to extract from PPTX files and create PowerPoint presentations"""
        try:
            if not self.b64_data or not self.b64_data.strip():
                return Message(text="❌ No base64 data provided")

            # Split comma-separated base64 data
            b64_parts = [part.strip() for part in self.b64_data.split(',')]
            b64_parts = [part for part in b64_parts if part]

            if not b64_parts:
                return Message(text="❌ No valid base64 data found after splitting")

            successful_contents = []
            errors = []
            total_powerpoints = 0

            # Process each PPTX file
            for i, b64_string in enumerate(b64_parts):
                message, success, content, count = self.process_single_pptx(b64_string, i + 1)

                if success:
                    successful_contents.append(content)
                    total_powerpoints += count
                else:
                    errors.append(message)

            # Build response
            if successful_contents:
                response_text = "🎯 COMBINED EXTRACTION & POWERPOINT CREATION RESULTS:\n"
                response_text += "=" * 70 + "\n\n"

                response_text += "\n".join(successful_contents)

                if errors:
                    response_text += "\n⚠️ Errors encountered:\n"
                    response_text += "\n".join(errors)

                response_text += f"\n--- FINAL SUMMARY ---\n"
                response_text += f"Files processed: {len(b64_parts)}\n"
                response_text += f"Successful files: {len(successful_contents)}\n"
                response_text += f"Total PowerPoints created: {total_powerpoints}\n"

            else:
                error_summary = "\n".join(errors) if errors else "Unknown errors occurred"
                response_text = f"❌ No valid PPTX files could be processed:\n\n{error_summary}"

            return Message(text=response_text)

        except Exception as e:
            error_text = f"❌ Failed to process PPTX files: {str(e)}"
            return Message(text=error_text)
