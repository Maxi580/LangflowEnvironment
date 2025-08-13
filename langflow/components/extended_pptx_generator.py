from langflow.custom import Component
from langflow.io import Output, MultilineInput
from langflow.schema import Message
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image
import tempfile
import base64
import os
import io
from dotenv import load_dotenv

load_dotenv()

PPTX_MAGIC_BYTES = os.getenv("PPTX_MAGIC_BYTES")
EXTENDED_TEMPLATE_PATH = os.getenv("EXTENDED_TEMPLATE_PATH")


class AtosExtendedTemplatePowerPointComponent(Component):
    display_name = "Atos Extended Template PowerPoint"
    description = "Creates PowerPoint from Atos extended template with Technology sections and additional fields"
    icon = "📊"
    inputs = [
        MultilineInput(
            name="customer_name",
            display_name="Customer Name",
            info="Customer name to replace {{CUSTOMER_NAME}} and search for logo",
            required=True,
        ),
        MultilineInput(
            name="about_client",
            display_name="About Client",
            info="Project name to replace {{ABOUT_CLIENT}}",
            required=True,
        ),
        MultilineInput(
            name="project_name",
            display_name="Project Name",
            info="Project name to replace {{PROJECT_NAME}}",
            required=True,
        ),
        MultilineInput(
            name="challenge_text",
            display_name="Challenge Description",
            info="Challenge description to replace {{CHALLENGE_TEXT}}",
            required=True,
        ),
        MultilineInput(
            name="solution_text",
            display_name="Solution Description",
            info="Solution description to replace {{SOLUTION_TEXT}}",
            required=True,
        ),
        MultilineInput(
            name="impact_text",
            display_name="Impact Description",
            info="Impact description to replace {{IMPACT_TEXT}}",
            required=True,
        ),
        MultilineInput(
            name="why_us_text",
            display_name="Why Us",
            info="Why Us description to replace {{WHY_US_TEXT}}",
            required=True,
        ),
        MultilineInput(
            name="it_impact_text",
            display_name="IT Impact",
            info="IT Impact description to replace {{IT_IMPACT}}",
            required=True,
        ),
        # New Technology Section Fields
        MultilineInput(
            name="methods_text",
            display_name="Methods Used",
            info="Methods and methodologies used to replace {{METHODS_TEXT}}",
            required=True,
        ),
        MultilineInput(
            name="software_used_impact",
            display_name="Software Tools and Systems",
            info="Software tools and systems used to replace {{SOFTWARE_USED_IMPACT}}",
            required=True,
        ),
        MultilineInput(
            name="hardware_text",
            display_name="Hardware Used",
            info="Hardware infrastructure used to replace {{HARDWARE_TEXT}}",
            required=True,
        ),
        MultilineInput(
            name="network_communication_text",
            display_name="Network and Communication Technology",
            info="Network and communication technology used to replace {{NETWORK_COMMUNICATION_TEXT}}",
            required=True,
        ),
        MultilineInput(
            name="technology_used_impact",
            display_name="Technology Used",
            info="Overall technology used to replace {{TECHNOLOGY_USED_IMPACT}}",
            required=True,
        ),
        MultilineInput(
            name="logo_base64",
            display_name="Company Logo (Base64)",
            info="Base64 encoded company logo image (optional - will replace placeholder logo)",
            required=False,
        ),
    ]

    outputs = [
        Output(
            display_name="Response",
            name="response_output",
            method="create_powerpoint"
        )
    ]

    def find_and_replace_text_in_slide(self, slide, replacements):
        """
        Find and replace text in all text boxes on a slide
        """
        replacements_made = 0

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text

                        for placeholder, replacement in replacements.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, replacement)
                                replacements_made += 1
                                print(f"✓ Replaced '{placeholder}' with '{replacement[:30]}...'")

            elif hasattr(shape, "table"):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for placeholder, replacement in replacements.items():
                                    if placeholder in run.text:
                                        run.text = run.text.replace(placeholder, replacement)
                                        replacements_made += 1
                                        print(f"✓ Replaced '{placeholder}' with '{replacement[:30]}...' (in table)")

        return replacements_made

    def decode_base64_image(self, base64_string):
        """
        Decode base64 image string and return image bytes
        """
        try:
            if base64_string.startswith('data:'):
                base64_string = base64_string.split(',', 1)[1]

            image_data = base64.b64decode(base64_string)

            img = Image.open(io.BytesIO(image_data))

            # Convert to RGB if necessary (for JPEG compatibility)
            if img.mode in ('RGBA', 'LA', 'P'):
                # Create white background
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                img = background

            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            return img_buffer.getvalue(), img.size

        except Exception as e:
            print(f"❌ Error decoding base64 image: {e}")
            return None, None

    def add_logo_at_fixed_position(self, slide, logo_data):
        try:
            cm_to_inches = 0.393701

            left = Inches(29.81 * cm_to_inches)
            top = Inches(0.81 * cm_to_inches)
            width = Inches(2.87 * cm_to_inches)
            height = Inches(2.53 * cm_to_inches)

            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_file.write(logo_data)
                tmp_logo_path = tmp_file.name

            try:
                pic = slide.shapes.add_picture(tmp_logo_path, left, top, width, height)

                print(f"✓ Added logo at position ({29.81:.2f}cm, {0.81:.2f}cm) with size {2.87:.2f}cm x {2.53:.2f}cm")
                return True

            finally:
                try:
                    os.unlink(tmp_logo_path)
                except:
                    pass

        except Exception as e:
            print(f"❌ Error adding logo: {e}")
            return False

    def add_logo_placeholder(self, slide):
        try:
            cm_to_inches = 0.393701

            left = Inches(29.81 * cm_to_inches)
            top = Inches(0.81 * cm_to_inches)
            width = Inches(2.87 * cm_to_inches)
            height = Inches(2.53 * cm_to_inches)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = "Logo Here"

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = 1
            font = paragraph.font
            font.size = Inches(0.2)
            font.color.rgb = RGBColor(128, 128, 128)

            line = textbox.line
            line.color.rgb = RGBColor(200, 200, 200)
            line.width = Inches(0.01)

            print(f"✓ Added logo placeholder at position ({29.81:.2f}cm, {0.81:.2f}cm)")
            return True

        except Exception as e:
            print(f"❌ Error adding logo placeholder: {e}")
            return False

    def create_powerpoint(self) -> Message:
        """Create PowerPoint from Atos extended template by replacing placeholders"""

        try:
            if not os.path.exists(EXTENDED_TEMPLATE_PATH):
                return (f"❌ Template file not found at {EXTENDED_TEMPLATE_PATH}. Please save your PowerPoint template as "
                        f"'template.pptx' in the same directory.")

            prs = Presentation(EXTENDED_TEMPLATE_PATH)
            print(f"✓ Loaded extended template with {len(prs.slides)} slides")

            # Extended replacements with all fields including technology sections
            replacements = {
                # Original fields
                '{{CUSTOMER_NAME}}': self.customer_name,
                '{{ABOUT_CLIENT}}': self.about_client,
                '{{PROJECT_NAME}}': self.project_name,
                '{{CHALLENGE_TEXT}}': self.challenge_text,
                '{{SOLUTION_TEXT}}': self.solution_text,
                '{{IMPACT_TEXT}}': self.impact_text,
                '{{WHY_US_TEXT}}': self.why_us_text,
                '{{IT_IMPACT}}': self.it_impact_text,
                # New Technology fields
                '{{METHODS_TEXT}}': self.methods_text,
                '{{SOFTWARE_USED_IMPACT}}': self.software_used_impact,
                '{{HARDWARE_TEXT}}': self.hardware_text,
                '{{NETWORK_COMMUNICATION_TEXT}}': self.network_communication_text,
                '{{TECHNOLOGY_USED_IMPACT}}': self.technology_used_impact,
            }

            has_logo = self.logo_base64 and self.logo_base64.strip()
            logo_data = None

            if has_logo:
                print("🖼️ Processing logo...")
                logo_data, logo_size = self.decode_base64_image(self.logo_base64.strip())
                if logo_data:
                    print(f"✓ Logo decoded successfully ({len(logo_data)} bytes)")
                else:
                    has_logo = False
                    print("❌ Failed to decode logo")

            # Process all slides
            total_replacements = 0
            for slide_idx, slide in enumerate(prs.slides):
                print(f"\n🔄 Processing slide {slide_idx + 1}...")

                # Replace text placeholders on this slide
                replacements_made = self.find_and_replace_text_in_slide(slide, replacements)
                total_replacements += replacements_made
                print(f"✓ Made {replacements_made} text replacements on slide {slide_idx + 1}")

                # Add logo only to the first slide (slide 0)
                if slide_idx == 0:
                    if has_logo and logo_data:
                        self.add_logo_at_fixed_position(slide, logo_data)
                    else:
                        self.add_logo_placeholder(slide)

            # Save the presentation to a temporary file
            temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            temp_file_path = temp_file.name
            temp_file.close()

            prs.save(temp_file_path)

            # Read the file content and encode it
            with open(temp_file_path, 'rb') as f:
                file_content = f.read()

            # Clean up the temporary file
            os.unlink(temp_file_path)

            # Encode to base64 for transmission
            base64_content = base64.b64encode(file_content).decode('utf-8')
            filename = f"tech_reference_{self.customer_name.replace(' ', '_').lower()}.pptx"

            # Return the response with the encoded file
            text_response = f"""Extended PowerPoint with Technology sections created successfully!

<{PPTX_MAGIC_BYTES}>
filename:{filename}
content_type:application/vnd.openxmlformats-officedocument.presentationml.presentation
size:{len(file_content)}
data:{base64_content}
</{PPTX_MAGIC_BYTES}>"""

            return text_response

        except Exception as e:
            return f"❌ Failed to create extended PowerPoint from template: {str(e)}"