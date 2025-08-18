from langflow.custom import Component
from langflow.io import Output, MultilineInput
from langflow.schema import Message
import base64
import tempfile
import os
from pptx import Presentation
from pptx.util import Cm
from PIL import Image
import io
from typing import List, Tuple


class PPTXReconstructorComponent(Component):
    display_name = "PPTX Reconstructor"
    description = "Reconstructs PPTX files from comma-separated base64 data and returns parsed content"
    icon = "📊"

    inputs = [
        MultilineInput(
            name="b64_data",
            display_name="Base64 PPTX Data",
            info="Comma-separated base64 encoded PPTX file data",
            required=True,
        ),
    ]

    outputs = [
        Output(
            display_name="Parsed Content",
            name="parsed_content_output",
            method="reconstruct_pptx_files"
        )
    ]

    def is_pptx_file(self, file_data: bytes) -> bool:
        """
        Check if the file data represents a valid PPTX file
        """
        try:
            if not file_data.startswith(b'PK\x03\x04'):
                return False

            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_file.write(file_data)
                temp_file_path = temp_file.name

            try:
                prs = Presentation(temp_file_path)
                _ = len(prs.slides)
                return True
            except Exception:
                return False
            finally:
                try:
                    os.unlink(temp_file_path)
                except:
                    pass

        except Exception:
            return False

    def extract_presentation_content(self, file_data: bytes) -> dict:
        """
        Extract and parse content from PPTX file using keyword-based title detection
        """
        try:
            with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as temp_file:
                temp_file.write(file_data)
                temp_file_path = temp_file.name

            try:
                prs = Presentation(temp_file_path)

                content = {
                    'slide_count': len(prs.slides),
                    'parsed_fields': {},
                    'logo_base64': '',
                    'company_name': '',
                    'project_name': ''
                }

                if len(prs.slides) >= 1:
                    company_name, project_name = self.extract_company_and_project_names(prs.slides[0])
                    content['company_name'] = company_name
                    content['project_name'] = project_name

                content['parsed_fields'] = self.parse_text_fields_from_slides(prs.slides)

                if len(prs.slides) >= 2:
                    content['logo_base64'] = self.extract_logo_from_slide(prs.slides[1])

                return content

            finally:
                try:
                    os.unlink(temp_file_path)
                except:
                    pass

        except Exception as e:
            return {'error': f'Failed to extract content: {str(e)}'}

    def extract_logo_from_slide(self, slide) -> str:
        """
        Extract logo image from slide using geo-based rectangle search
        """
        try:
            target_left = Cm(0)
            target_top = Cm(1)
            target_width = Cm(6.52)
            target_height = Cm(3.04)

            target_left_emu = target_left
            target_top_emu = target_top
            target_right_emu = target_left + target_width
            target_bottom_emu = target_top + target_height

            for shape in slide.shapes:
                if hasattr(shape, 'image'):
                    shape_right = shape.left + shape.width
                    shape_bottom = shape.top + shape.height

                    overlaps = (shape.left < target_right_emu and
                                shape_right > target_left_emu and
                                shape.top < target_bottom_emu and
                                shape_bottom > target_top_emu)

                    if overlaps:
                        image_stream = shape.image.blob
                        image = Image.open(io.BytesIO(image_stream))

                        img_buffer = io.BytesIO()
                        image.save(img_buffer, format='PNG')
                        img_buffer.seek(0)

                        base64_image = base64.b64encode(img_buffer.getvalue()).decode('utf-8')
                        return base64_image

            return ""

        except Exception as e:
            return ""

    def extract_company_and_project_names(self, slide) -> tuple:
        """
        Extract company name and project name from slide 1 by finding the two highest text fields
        """
        try:
            text_shapes = self.get_text_shapes_from_slide(slide)

            if len(text_shapes) < 2:
                return "", ""

            # Filter to text blocks up to 150 characters as potential titles
            # Also exclude common non-title text
            title_candidates = []
            exclude_phrases = [
                'internal use only',
                'confidential',
                'draft',
                'do not distribute',
                'proprietary'
            ]

            for text_shape in text_shapes:
                text_lower = text_shape['text'].lower().strip()

                # Skip if text is too long
                if len(text_shape['text']) > 150:
                    continue

                # Skip if text contains excluded phrases
                should_exclude = False
                for exclude_phrase in exclude_phrases:
                    if exclude_phrase in text_lower:
                        should_exclude = True
                        break

                if not should_exclude:
                    title_candidates.append(text_shape)

            if len(title_candidates) < 2:
                # If we don't have enough candidates after filtering, fall back to all text shapes
                # but still exclude the specific phrases
                title_candidates = []
                for text_shape in text_shapes:
                    text_lower = text_shape['text'].lower().strip()
                    should_exclude = False
                    for exclude_phrase in exclude_phrases:
                        if exclude_phrase in text_lower:
                            should_exclude = True
                            break
                    if not should_exclude:
                        title_candidates.append(text_shape)

            if len(title_candidates) < 1:
                return "", ""

            # Sort by Y position (top to bottom)
            title_candidates.sort(key=lambda x: x['top'])

            company_name = title_candidates[0]['text'].strip()
            project_name = title_candidates[1]['text'].strip() if len(title_candidates) > 1 else ""

            return company_name, project_name

        except Exception as e:
            return "", ""

    def parse_text_fields_from_slides(self, slides) -> dict:
        """
        Parse text fields from all slides using keyword patterns
        """
        keyword_patterns = {
            'the_client': ['the client', 'client'],
            'the_challenge': ['the challenge', 'challenge'],
            'solution': ['solution'],
            'impact': ['impact'],
            'why_eviden': ['why eviden', 'eviden']
        }

        all_parsed_fields = {}

        for slide_idx, slide in enumerate(slides):
            text_shapes = self.get_text_shapes_from_slide(slide)
            slide_fields = self.extract_content_by_keywords(text_shapes, keyword_patterns)

            for field_name, content in slide_fields.items():
                if content.strip():
                    all_parsed_fields[field_name] = content

        return all_parsed_fields

    def get_text_shapes_from_slide(self, slide) -> list:
        """
        Extract all text shapes from a slide with their spatial information
        """
        text_shapes = []

        for shape in slide.shapes:
            if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
                text = shape.text_frame.text.strip()
                text_shapes.append({
                    'text': text,
                    'text_length': len(text),
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })

        return text_shapes

    def extract_content_by_keywords(self, text_shapes, keyword_patterns) -> dict:
        """
        Extract content based on keyword patterns using spatial logic
        """
        result = {}
        title_positions = {}

        for category, keywords in keyword_patterns.items():
            for text_shape in text_shapes:
                text = text_shape['text'].strip().lower()

                if len(text) > 25:
                    continue

                for keyword in keywords:
                    if keyword in text:
                        title_positions[category] = {
                            'x': text_shape['left'],
                            'y': text_shape['top'],
                            'text': text_shape['text'],
                            'shape_info': text_shape
                        }
                        break

                if category in title_positions:
                    break

        for category, title_info in title_positions.items():
            content = self.find_content_below_title(text_shapes, title_info, title_positions)
            result[category] = content

        return result

    def find_content_below_title(self, text_shapes, title_info, all_title_positions) -> str:
        """
        Find the first text field below a title with spatial proximity rules
        """
        title_x = title_info['x']
        title_y = title_info['y']

        candidate_content = []
        x_tolerance = 2 * 360000  # 2cm tolerance

        for text_shape in text_shapes:
            if (text_shape['left'] == title_x and text_shape['top'] == title_y):
                continue

            is_other_title = False
            for other_category, other_title_info in all_title_positions.items():
                if (text_shape['left'] == other_title_info['x'] and
                        text_shape['top'] == other_title_info['y']):
                    is_other_title = True
                    break

            if is_other_title:
                continue

            if text_shape['top'] > title_y:
                x_distance = abs(text_shape['left'] - title_x)

                if x_distance < x_tolerance:
                    y_distance = text_shape['top'] - title_y
                    candidate_content.append({
                        'text': text_shape['text'],
                        'x_distance': x_distance,
                        'y_distance': y_distance,
                        'text_length': text_shape['text_length']
                    })

        if candidate_content:
            candidate_content.sort(key=lambda x: (x['y_distance'], -x['text_length']))
            selected_content = candidate_content[0]
            return selected_content['text']
        else:
            return ""

    def clean_b64_string(self, b64_string: str) -> str:
        """
        Clean base64 string by removing data URI prefixes and whitespace
        """
        b64_string = b64_string.strip()

        if b64_string.startswith('data:'):
            if ',' in b64_string:
                b64_string = b64_string.split(',', 1)[1]

        b64_string = ''.join(b64_string.split())

        return b64_string

    def process_single_pptx(self, b64_string: str, index: int) -> Tuple[str, bool, dict]:
        """
        Process a single PPTX file from base64 string
        """
        try:
            clean_b64 = self.clean_b64_string(b64_string)

            if not clean_b64:
                return f"⚠️ Empty base64 string for file {index}", False, {}

            file_data = base64.b64decode(clean_b64)

            if not self.is_pptx_file(file_data):
                return f"❌ File {index} is not a valid PPTX format", False, {}

            content = self.extract_presentation_content(file_data)

            return f"✅ Successfully processed PPTX file {index}", True, content

        except Exception as e:
            error_msg = f"❌ Error processing file {index}: {str(e)}"
            return error_msg, False, {}

    def format_content_output(self, content: dict) -> str:
        """
        Format the extracted and parsed content into a structured format
        """
        if 'error' in content:
            return f"❌ Error extracting content: {content['error']}"

        output_lines = []

        if 'parsed_fields' in content:
            parsed_fields = content['parsed_fields']

            # Add company and project names first
            company_name = content.get('company_name', '')
            project_name = content.get('project_name', '')

            output_lines.append("Customer Name:")
            if company_name:
                output_lines.append(company_name)
            else:
                output_lines.append("[No company name found]")
            output_lines.append("")

            output_lines.append("Project Name:")
            if project_name:
                output_lines.append(project_name)
            else:
                output_lines.append("[No project name found]")
            output_lines.append("")

            # Define the order and proper display names for other fields
            field_mapping = {
                'the_client': 'About Client',
                'the_challenge': 'Challenge',
                'solution': 'Solution',
                'impact': 'Impact',
                'why_eviden': 'Why Us'
            }

            for field_key, display_name in field_mapping.items():
                field_content = parsed_fields.get(field_key, '')

                output_lines.append(f"{display_name}:")
                if field_content:
                    output_lines.append(field_content)
                else:
                    output_lines.append("[No content found]")
                output_lines.append("")

            # Add logo section
            logo_base64 = content.get('logo_base64', '')
            output_lines.append("Logo:")
            if logo_base64:
                output_lines.append(logo_base64)
            else:
                output_lines.append("[No logo found]")
            output_lines.append("")

            # Clean end marker for easy filtering
            output_lines.append("--- END OF EXTRACTION ---")
        else:
            output_lines.append("❌ No parsed fields found")

        return "\n".join(output_lines)

    def reconstruct_pptx_files(self) -> Message:
        """
        Main method to process PPTX files from comma-separated base64 data
        """
        try:
            if not self.b64_data or not self.b64_data.strip():
                return Message(text="❌ No base64 data provided")

            b64_parts = [part.strip() for part in self.b64_data.split(',')]
            b64_parts = [part for part in b64_parts if part]

            if not b64_parts:
                return Message(text="❌ No valid base64 data found after splitting")

            successful_contents = []
            errors = []

            for i, b64_string in enumerate(b64_parts):
                message, success, content = self.process_single_pptx(b64_string, i + 1)

                if success:
                    successful_contents.append(content)
                else:
                    errors.append(message)

            if successful_contents:
                response_lines = []

                if len(successful_contents) == 1:
                    response_lines.append(self.format_content_output(successful_contents[0]))
                else:
                    for i, content in enumerate(successful_contents):
                        response_lines.append(f"📊 File {i + 1}:")
                        response_lines.append(self.format_content_output(content))
                        response_lines.append("-" * 50)

                if errors:
                    response_lines.extend([
                        "",
                        "⚠️ Errors encountered:",
                        *errors
                    ])

                response_text = "\n".join(response_lines)

            else:
                error_summary = "\n".join(errors) if errors else "Unknown errors occurred"
                response_text = f"❌ No valid PPTX files could be processed:\n\n{error_summary}"

            return Message(text=response_text)

        except Exception as e:
            error_text = f"❌ Failed to process PPTX files: {str(e)}"
            return Message(text=error_text)
