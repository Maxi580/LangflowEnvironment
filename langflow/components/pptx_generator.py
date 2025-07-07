from langflow.custom import Component
from langflow.io import Output, MultilineInput
from langflow.schema import Message
from pptx import Presentation
from pptx.util import Inches
from PIL import Image
import tempfile
import base64
import os
import io
from dotenv import load_dotenv

load_dotenv()

PPTX_MAGIC_BYTES = os.getenv("PPTX_MAGIC_BYTES")
TEMPLATE_PATH = os.getenv("TEMPLATE_PATH")


class AtosTemplatePowerPointComponent(Component):
    display_name = "Atos Template PowerPoint"
    description = "Creates PowerPoint from Atos template by replacing placeholder text and adding logo"
    icon = "📊"
    inputs = [
        MultilineInput(
            name="customer_name",
            display_name="Customer Name",
            info="Customer name to replace {{CUSTOMER_NAME}} and search for logo",
            required=True,
        ),
        MultilineInput(
            name="about_client",
            display_name="About Client",
            info="Project name to replace {{ABOUT_CLIENT}}",
            required=True,
        ),
        MultilineInput(
            name="project_name",
            display_name="Project Name",
            info="Project name to replace {{PROJECT_NAME}}",
            required=True,
        ),
        MultilineInput(
            name="challenge_text",
            display_name="Challenge Description",
            info="Challenge description to replace {{CHALLENGE_TEXT}}",
            required=True,
        ),
        MultilineInput(
            name="solution_text",
            display_name="Solution Description",
            info="Solution description to replace {{SOLUTION_TEXT}}",
            required=True,
        ),
        MultilineInput(
            name="impact_text",
            display_name="Impact Description",
            info="Impact description to replace {{IMPACT_TEXT}}",
            required=True,
        ),
        MultilineInput(
            name="logo_base64",
            display_name="Company Logo (Base64)",
            info="Base64 encoded company logo image (optional - will replace placeholder logo)",
            required=False,
        ),
    ]

    outputs = [
        Output(
            display_name="Response",
            name="response_output",
            method="create_powerpoint"
        )
    ]

    def find_and_replace_text_in_slide(self, slide, replacements):
        """
        Find and replace text in all text boxes on a slide
        """
        replacements_made = 0

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text

                        for placeholder, replacement in replacements.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, replacement)
                                replacements_made += 1
                                print(f"✓ Replaced '{placeholder}' with '{replacement[:30]}...'")

            elif hasattr(shape, "table"):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for placeholder, replacement in replacements.items():
                                    if placeholder in run.text:
                                        run.text = run.text.replace(placeholder, replacement)
                                        replacements_made += 1
                                        print(f"✓ Replaced '{placeholder}' with '{replacement[:30]}...' (in table)")

        return replacements_made

    def decode_base64_image(self, base64_string):
        """
        Decode base64 image string and return image bytes
        """
        try:
            if base64_string.startswith('data:'):
                base64_string = base64_string.split(',', 1)[1]

            image_data = base64.b64decode(base64_string)

            img = Image.open(io.BytesIO(image_data))

            # Convert to RGB if necessary (for JPEG compatibility)
            if img.mode in ('RGBA', 'LA', 'P'):
                # Create white background
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                img = background

            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            return img_buffer.getvalue(), img.size

        except Exception as e:
            print(f"❌ Error decoding base64 image: {e}")
            return None, None

    def add_logo_at_fixed_position(self, slide, logo_data):
        try:
            cm_to_inches = 0.393701

            left = Inches(29.81 * cm_to_inches)
            top = Inches(0.81 * cm_to_inches)
            width = Inches(2.87 * cm_to_inches)
            height = Inches(2.53 * cm_to_inches)

            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_file.write(logo_data)
                tmp_logo_path = tmp_file.name

            try:
                pic = slide.shapes.add_picture(tmp_logo_path, left, top, width, height)

                print(f"✓ Added logo at position ({29.81:.2f}cm, {0.81:.2f}cm) with size {2.87:.2f}cm x {2.53:.2f}cm")
                return True

            finally:
                try:
                    os.unlink(tmp_logo_path)
                except:
                    pass

        except Exception as e:
            print(f"❌ Error adding logo: {e}")
            return False

    def create_powerpoint(self) -> Message:
        """Create PowerPoint from Atos template by replacing placeholders"""

        try:
            if not os.path.exists(TEMPLATE_PATH):
                return (f"❌ Template file not found at {TEMPLATE_PATH}. Please save your PowerPoint template as "
                        f"'template.pptx' in the same directory.")

            prs = Presentation(TEMPLATE_PATH)
            print(f"✓ Loaded template with {len(prs.slides)} slides")

            replacements = {
                '{{CUSTOMER_NAME}}': self.customer_name,
                '{{ABOUT_CLIENT}}': self.about_client,
                '{{PROJECT_NAME}}': self.project_name,
                '{{CHALLENGE_TEXT}}': self.challenge_text,
                '{{SOLUTION_TEXT}}': self.solution_text,
                '{{IMPACT_TEXT}}': self.impact_text,
            }

            logo_data = None
            logo_size = None
            if self.logo_base64 and self.logo_base64.strip():
                print("🖼️ Processing logo...")
                logo_data, logo_size = self.decode_base64_image(self.logo_base64.strip())
                if logo_data:
                    print(f"✓ Logo decoded successfully ({len(logo_data)} bytes)")
                else:
                    print("❌ Failed to decode logo")

            total_replacements = 0
            logo_replacements = 0

            for slide_idx, slide in enumerate(prs.slides):
                print(f"\n🔄 Processing slide {slide_idx + 1}...")

                replacements_made = self.find_and_replace_text_in_slide(slide, replacements)
                total_replacements += replacements_made

                if logo_data and slide_idx == 0:
                    if self.add_logo_at_fixed_position(slide, logo_data):
                        logo_replacements += 1

            print(f"\n✅ Made {total_replacements} text replacements")
            if logo_data:
                print(f"✅ Made {logo_replacements} logo replacements")

            temp_file = tempfile.NamedTemporaryFile(suffix='.pptx', delete=False)
            temp_file_path = temp_file.name
            temp_file.close()

            prs.save(temp_file_path)

            with open(temp_file_path, 'rb') as f:
                file_content = f.read()

            os.unlink(temp_file_path)

            base64_content = base64.b64encode(file_content).decode('utf-8')
            filename = f"reference_{self.customer_name.replace(' ', '_').lower()}.pptx"

            text_response = f"""PowerPoint created successfully!

<{PPTX_MAGIC_BYTES}>
filename:{filename}
content_type:application/vnd.openxmlformats-officedocument.presentationml.presentation
size:{len(file_content)}
data:{base64_content}
</{PPTX_MAGIC_BYTES}>"""

            return text_response

        except Exception as e:
            return f"❌ Failed to create PowerPoint from template: {str(e)}"
