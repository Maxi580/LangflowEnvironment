import json
import re
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image
import tempfile
import base64
import os
import io
from dotenv import load_dotenv

from langflow.custom import Component
from langflow.io import MultilineInput, HandleInput, Output
from langflow.schema.message import Message

load_dotenv()

PPTX_MAGIC_BYTES = os.getenv("PPTX_MAGIC_BYTES")
TEMPLATE_PATH = os.getenv("TEMPLATE_PATH")


class MultiLanguagePowerPointComponent(Component):
    display_name = "Multi-Language PowerPoint"
    description = (
        "Takes the Translation Engine JSON and creates one PowerPoint per language. "
        "Each PPTX uses the Atos template with translated sections filled in."
    )
    icon = "📊"

    inputs = [
        HandleInput(
            name="translations_json",
            display_name="Translations JSON",
            info="The JSON output from the Translation Engine.",
            input_types=["Message"],
            is_list=False,
        ),
        MultilineInput(
            name="logo_base64",
            display_name="Company Logo (Base64)",
            info="Base64 encoded company logo image (optional).",
            required=False,
        ),
        MultilineInput(
            name="create_pptx_flag",
            display_name="Create PPTX Flag",
            info="'true' or 'false' — skip generation if false.",
            required=False,
            value="true",
        ),
    ]

    outputs = [
        Output(
            display_name="Response",
            name="response_output",
            method="create_powerpoints",
        ),
    ]

    # ------------------------------------------------------------------ #
    #  Section → placeholder mapping                                       #
    # ------------------------------------------------------------------ #

    # Maps Translation Engine section keys to PPTX template placeholders
    SECTION_TO_PLACEHOLDER = {
        "Client_Name":      "{{CUSTOMER_NAME}}",
        "About_Client":     "{{ABOUT_CLIENT}}",
        "Project":          "{{PROJECT_NAME}}",
        "Challenge":        "{{CHALLENGE_TEXT}}",
        "Solution":         "{{SOLUTION_TEXT}}",
        "Business_Impact":  "{{IMPACT_TEXT}}",
    }

    # ------------------------------------------------------------------ #
    #  PPTX helpers (from original component)                              #
    # ------------------------------------------------------------------ #

    def _find_and_replace_text_in_slide(self, slide, replacements):
        replacements_made = 0
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for placeholder, replacement in replacements.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, replacement)
                                replacements_made += 1

            elif hasattr(shape, "table"):
                for row in shape.table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for placeholder, replacement in replacements.items():
                                    if placeholder in run.text:
                                        run.text = run.text.replace(placeholder, replacement)
                                        replacements_made += 1

        return replacements_made

    def _decode_base64_image(self, base64_string):
        try:
            if base64_string.startswith("data:"):
                base64_string = base64_string.split(",", 1)[1]

            image_data = base64.b64decode(base64_string)
            img = Image.open(io.BytesIO(image_data))

            if img.mode in ("RGBA", "LA", "P"):
                background = Image.new("RGB", img.size, (255, 255, 255))
                if img.mode == "P":
                    img = img.convert("RGBA")
                background.paste(img, mask=img.split()[-1] if img.mode in ("RGBA", "LA") else None)
                img = background

            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            return buf.getvalue(), img.size

        except Exception as e:
            print(f"Error decoding base64 image: {e}")
            return None, None

    def _add_logo_at_fixed_position(self, slide, logo_data):
        try:
            cm = 0.393701
            left   = Inches(29.81 * cm)
            top    = Inches(0.81  * cm)
            width  = Inches(2.87  * cm)
            height = Inches(2.53  * cm)

            with tempfile.NamedTemporaryFile(suffix=".png", delete=False) as tmp:
                tmp.write(logo_data)
                tmp_path = tmp.name

            try:
                slide.shapes.add_picture(tmp_path, left, top, width, height)
                return True
            finally:
                try:
                    os.unlink(tmp_path)
                except Exception:
                    pass

        except Exception as e:
            print(f"Error adding logo: {e}")
            return False

    def _add_logo_placeholder(self, slide):
        try:
            cm = 0.393701
            left   = Inches(29.81 * cm)
            top    = Inches(0.81  * cm)
            width  = Inches(2.87  * cm)
            height = Inches(2.53  * cm)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            tf = textbox.text_frame
            tf.text = "Logo Here"
            p = tf.paragraphs[0]
            p.alignment = 1
            p.font.size = Inches(0.2)
            p.font.color.rgb = RGBColor(128, 128, 128)
            textbox.line.color.rgb = RGBColor(200, 200, 200)
            textbox.line.width = Inches(0.01)
            return True

        except Exception as e:
            print(f"Error adding logo placeholder: {e}")
            return False

    # ------------------------------------------------------------------ #
    #  Build one PPTX for a single language                                #
    # ------------------------------------------------------------------ #

    def _create_single_pptx(self, language: str, sections: dict, logo_data=None) -> tuple[bytes, str] | None:
        """
        Create one PPTX for the given language.
        Returns (file_bytes, filename) or None on failure.
        """
        if not os.path.exists(TEMPLATE_PATH):
            print(f"Template not found: {TEMPLATE_PATH}")
            return None

        prs = Presentation(TEMPLATE_PATH)

        # Build replacements from sections
        replacements = {}
        for section_key, placeholder in self.SECTION_TO_PLACEHOLDER.items():
            value = sections.get(section_key, "")
            if value:
                replacements[placeholder] = value

        # Process slides
        for slide_idx, slide in enumerate(prs.slides):
            self._find_and_replace_text_in_slide(slide, replacements)

            # Logo on first slide
            if slide_idx == 0:
                if logo_data:
                    self._add_logo_at_fixed_position(slide, logo_data)
                else:
                    self._add_logo_placeholder(slide)

        # Save to bytes
        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        tmp_path = tmp.name
        tmp.close()

        prs.save(tmp_path)
        with open(tmp_path, "rb") as f:
            file_bytes = f.read()
        os.unlink(tmp_path)

        # Build filename
        customer = sections.get("Client_Name", "client").replace(" ", "_").lower()
        lang_slug = language.replace(" ", "_").lower()
        filename = f"reference_{customer}_{lang_slug}.pptx"

        return file_bytes, filename

    # ------------------------------------------------------------------ #
    #  Main output                                                         #
    # ------------------------------------------------------------------ #

    def create_powerpoints(self) -> Message:
        # Check flag
        flag = (self.create_pptx_flag or "true").strip().lower()
        if flag == "false":
            print("create_pptx is false — skipping PowerPoint generation.")
            return Message(text="")

        # Parse translations JSON
        raw = ""
        if isinstance(self.translations_json, Message):
            raw = self.translations_json.text or ""
        elif isinstance(self.translations_json, str):
            raw = self.translations_json
        else:
            raw = str(self.translations_json)

        if not raw.strip():
            return Message(text="No translation data received.")

        try:
            data = json.loads(raw)
        except json.JSONDecodeError as e:
            return Message(text=f"Failed to parse translations JSON: {e}")

        translations = data.get("translations", [])
        if not translations:
            return Message(text="No translations found in JSON.")

        # Decode logo once
        logo_data = None
        if self.logo_base64 and self.logo_base64.strip():
            logo_data, _ = self._decode_base64_image(self.logo_base64.strip())

        # Generate one PPTX per language
        response_parts = []
        files_created = 0

        for entry in translations:
            language = entry.get("language", "Unknown")
            sections = entry.get("sections", {})
            error = entry.get("error")

            if error:
                response_parts.append(f"⚠️ Skipped {language}: {error}")
                continue

            if not sections:
                response_parts.append(f"⚠️ Skipped {language}: no sections.")
                continue

            result = self._create_single_pptx(language, sections, logo_data)
            if result is None:
                response_parts.append(f"❌ Failed to create PPTX for {language}.")
                continue

            file_bytes, filename = result
            b64_content = base64.b64encode(file_bytes).decode("utf-8")

            response_parts.append(
                f"{language}: {filename}\n\n"
                f"<{PPTX_MAGIC_BYTES}>\n"
                f"filename:{filename}\n"
                f"content_type:application/vnd.openxmlformats-officedocument.presentationml.presentation\n"
                f"size:{len(file_bytes)}\n"
                f"data:{b64_content}\n"
                f"</{PPTX_MAGIC_BYTES}>"
            )
            files_created += 1

        self.status = f"Created {files_created} PPTX file(s) for {len(translations)} language(s)."

        return Message(text="\n\n".join(response_parts))