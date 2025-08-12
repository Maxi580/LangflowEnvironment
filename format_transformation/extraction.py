import os
from pptx import Presentation
from pptx.util import Cm
from PIL import Image
import io
import re


def extract_text_and_image_from_slide(ppt_path, output_dir="extracted_content"):
    """
    Extract image and structured text from slides 1 and 2 in a PowerPoint presentation.

    Args:
        ppt_path (str): Path to the PowerPoint file
        output_dir (str): Directory to save extracted content

    Returns:
        dict: Dictionary containing extracted image path and text content, or None if extraction failed
    """

    os.makedirs(output_dir, exist_ok=True)

    try:
        # Load the presentation
        prs = Presentation(ppt_path)

        # Check if we have at least 2 slides
        if len(prs.slides) < 2:
            print("Error: Presentation doesn't have at least 2 slides")
            return None

        # Extract image from slide 2 (index 1) - keeping existing logic
        slide2 = prs.slides[1]
        image_info = extract_image_from_slide_helper(slide2, ppt_path, output_dir)

        text_content = {}

        # Extract from slide 1 (for "The client" and "The challenge")
        slide1 = prs.slides[0]
        slide1_content = extract_structured_text(slide1, image_info, slide_number=1)
        text_content.update(slide1_content)

        # Extract from slide 2 (for existing fields)
        slide2_content = extract_structured_text(slide2, image_info, slide_number=2)
        text_content.update(slide2_content)

        result = {
            'image_path': image_info['path'] if image_info else None,
            'text_content': text_content,
            'file_name': os.path.basename(ppt_path)
        }

        return result

    except Exception as e:
        print(f"Error processing PowerPoint file: {str(e)}")
        return None


def extract_image_from_slide_helper(slide, ppt_path, output_dir):
    """Helper function to extract image from slide"""

    # Define the target area for image extraction
    target_left = Cm(0)
    target_top = Cm(1)
    target_width = Cm(6.52)
    target_height = Cm(3.04)

    # Convert to EMUs for comparison
    target_left_emu = target_left
    target_top_emu = target_top
    target_right_emu = target_left + target_width
    target_bottom_emu = target_top + target_height

    # Iterate through all shapes in the slide
    for shape in slide.shapes:
        # Check if shape is a picture
        if hasattr(shape, 'image'):
            # Check if the image is within our target area
            shape_right = shape.left + shape.width
            shape_bottom = shape.top + shape.height

            # Check for overlap with target area
            overlaps = (shape.left < target_right_emu and
                        shape_right > target_left_emu and
                        shape.top < target_bottom_emu and
                        shape_bottom > target_top_emu)

            if overlaps:
                # Extract the image
                image_stream = shape.image.blob
                image = Image.open(io.BytesIO(image_stream))

                # Generate output filename
                base_name = os.path.splitext(os.path.basename(ppt_path))[0]
                output_filename = f"{base_name}_slide2_image.png"
                output_path = os.path.join(output_dir, output_filename)

                # Save the image
                image.save(output_path, "PNG")
                print(f"Image extracted and saved to: {output_path}")
                return {
                    'path': output_path,
                    'bottom_y': shape_bottom
                }

    print("No image found in the specified area")
    return None


def extract_structured_text(slide, image_info, slide_number=None):
    """
    Extract structured text from slide by finding titles and their corresponding content.
    Enhanced version with slide 1 and slide 2 specific title patterns.
    """

    title_patterns = {
        'the_client': [
            r'^The client$',
            r'^The Client$',
            r'^Client$',
            r'^THE CLIENT$'
        ],
        'the_challenge': [
            r'^The challenge$',
            r'^The Challenge$',
            r'^Challenge$',
            r'^THE CHALLENGE$'
        ],
        'solution': [
            r'^The solution$',
            r'^The Solution$',
            r'^Solution$'
        ],
        'impact': [
            r'^The impact$',
            r'^The Impact$',
            r'^Impact$'
        ],
        'why_eviden': [
            r'^Why Eviden\??$',
            r'^Why Eviden$',
            r'^Eviden$'
        ]
    }

    # Step 1: Find all text shapes and categorize them
    text_shapes = []
    for shape in slide.shapes:
        if hasattr(shape, 'text_frame') and shape.text_frame.text.strip():
            text = shape.text_frame.text.strip()
            text_shapes.append({
                'shape': shape,
                'text': text,
                'text_length': len(text),
                'left': shape.left,
                'top': shape.top,
                'width': shape.width,
                'height': shape.height
            })

    print(f"Found {len(text_shapes)} text shapes on slide {slide_number if slide_number else 'unknown'}")

    # Debug: Print all text shapes with their properties
    for i, text_shape in enumerate(text_shapes):
        print(f"Shape {i}: x={text_shape['left']}, y={text_shape['top']}, "
              f"len={text_shape['text_length']}, text='{text_shape['text'][:50]}...'")

    # Step 2: Find titles using stricter matching
    title_positions = {}

    for category, patterns in title_patterns.items():
        for text_shape in text_shapes:
            text = text_shape['text'].strip()

            # Only consider short text blocks as potential titles (less than 50 characters)
            if len(text) > 25:
                continue

            for pattern in patterns:
                if re.match(pattern, text, re.IGNORECASE):
                    title_positions[category] = {
                        'x': text_shape['left'],
                        'y': text_shape['top'],
                        'text': text_shape['text'],
                        'shape_info': text_shape
                    }
                    print(f"Found {category} title: '{text}' at x={text_shape['left']}, y={text_shape['top']}")
                    break
            if category in title_positions:
                break

    if not title_positions:
        print("No titles found with strict matching, trying flexible matching...")
        title_positions = find_titles_flexible(text_shapes, title_patterns)

    company_name, project_name = extract_company_and_project_names(
        text_shapes, image_info, title_positions
    )

    result = {}
    if title_positions:
        result = extract_content_by_titles(text_shapes, title_positions)
    else:
        print("No titles found")

    result['company_name'] = company_name
    result['project_name'] = project_name

    return result


def extract_company_and_project_names(text_shapes, image_info, title_positions):
    """
    Extract company name and project name based on spatial positioning.
    Company name should be higher (lower y value) than project name.
    Both should be between the image bottom and the highest title.
    """

    if not image_info:
        print("No image info available for spatial extraction")
        return "", ""

    image_bottom_y = image_info['bottom_y']

    # Find the highest (lowest y value) title position
    if not title_positions:
        print("No title positions available for spatial extraction")
        return "", ""

    highest_title_y = min(title_info['y'] for title_info in title_positions.values())

    print(f"Image bottom Y: {image_bottom_y}")
    print(f"Highest title Y: {highest_title_y}")

    # Find text shapes that are between image bottom and highest title
    candidate_shapes = []

    for text_shape in text_shapes:
        # Check if shape is between image and titles vertically
        if image_bottom_y < text_shape['top'] < highest_title_y:
            # Skip if this is actually one of our identified titles
            is_title = False
            for title_info in title_positions.values():
                if (text_shape['left'] == title_info['x'] and
                        text_shape['top'] == title_info['y']):
                    is_title = True
                    break

            if not is_title:
                candidate_shapes.append(text_shape)
                print(f"Candidate shape: '{text_shape['text'][:50]}...' at y={text_shape['top']}")

    if len(candidate_shapes) < 2:
        print(f"Warning: Found only {len(candidate_shapes)} candidate shapes for company/project names")
        # Handle case where we have fewer than 2 candidates
        if len(candidate_shapes) == 1:
            # Assume single field is company name
            company_name = candidate_shapes[0]['text']
            project_name = ""
        else:
            company_name = ""
            project_name = ""
    else:
        # Sort by Y position (top to bottom)
        candidate_shapes.sort(key=lambda x: x['top'])

        # First (highest) is company name, second is project name
        company_name = candidate_shapes[0]['text']
        project_name = candidate_shapes[1]['text']

        print(f"Extracted company name: '{company_name}'")
        print(f"Extracted project name: '{project_name}'")

    return company_name, project_name


def find_titles_flexible(text_shapes, title_patterns):
    """
    Flexible title finding - looks for title keywords in short text blocks
    Enhanced with slide 1 keywords
    """
    title_positions = {}

    # Create simpler keyword patterns - now includes slide 1 keywords
    keyword_patterns = {
        'the_client': ['the client', 'client'],
        'the_challenge': ['the challenge', 'challenge'],
        'solution': ['solution'],
        'impact': ['impact'],
        'why_eviden': ['why eviden', 'eviden']
    }

    for category, keywords in keyword_patterns.items():
        for text_shape in text_shapes:
            text = text_shape['text'].strip().lower()

            # Only consider relatively short text blocks as potential titles
            if len(text) > 25:
                continue

            for keyword in keywords:
                if keyword in text:
                    title_positions[category] = {
                        'x': text_shape['left'],
                        'y': text_shape['top'],
                        'text': text_shape['text'],
                        'shape_info': text_shape
                    }
                    print(f"Found {category} title (flexible): '{text_shape['text'][:50]}...' "
                          f"at x={text_shape['left']}, y={text_shape['top']}")
                    break
            if category in title_positions:
                break

    return title_positions


def extract_content_by_titles(text_shapes, title_positions):
    """
    Extract content based on found titles using improved spatial logic
    """
    result = {}

    for category, title_info in title_positions.items():
        title_x = title_info['x']
        title_y = title_info['y']

        print(f"\nLooking for content for {category} title at x={title_x}, y={title_y}")

        candidate_content = []

        for text_shape in text_shapes:
            # Skip the title itself
            if (text_shape['left'] == title_x and
                    text_shape['top'] == title_y):
                continue

            # Skip other titles
            is_other_title = False
            for other_category, other_title_info in title_positions.items():
                if (other_category != category and
                        text_shape['left'] == other_title_info['x'] and
                        text_shape['top'] == other_title_info['y']):
                    is_other_title = True
                    break

            if is_other_title:
                continue

            # Check if this content could belong to this title
            # Content should be below the title (higher y value)
            if text_shape['top'] > title_y:
                # Calculate horizontal proximity
                x_distance = abs(text_shape['left'] - title_x)
                x_tolerance = Cm(2)

                if x_distance < x_tolerance:
                    y_distance = text_shape['top'] - title_y
                    candidate_content.append({
                        'text': text_shape['text'],
                        'x_distance': x_distance,
                        'y_distance': y_distance,
                        'text_length': len(text_shape['text']),
                        'shape_info': text_shape
                    })

        # If there are several candidates that have (almost) the same x and higher y than title, then pick lowest y
        if candidate_content:
            # Sort by vertical distance (closest to title first) then by text length
            candidate_content.sort(key=lambda x: (x['y_distance'], -x['text_length']))

            selected_content = candidate_content[0]
            result[category] = selected_content['text']
            print(f"{category}: Selected content with {selected_content['text_length']} characters")
            print(f"  Content preview: {selected_content['text'][:100]}...")
        else:
            result[category] = ""
            print(f"{category}: No content blocks found below title")

    return result
