from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image
import tempfile
import base64
import os
import io
import argparse

# Configuration
TEMPLATE_PATH = os.getenv("TEMPLATE_PATH", "template.pptx")


class PowerPointGenerator:
    def __init__(self, template_path=None):
        """Initialize the PowerPoint generator with template path"""
        self.template_path = template_path or TEMPLATE_PATH

    def find_and_replace_text_in_slide(self, slide, replacements):
        """
        Find and replace text in all text boxes on a slide
        """
        replacements_made = 0

        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                text_frame = shape.text_frame

                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        original_text = run.text

                        for placeholder, replacement in replacements.items():
                            if placeholder in run.text:
                                run.text = run.text.replace(placeholder, replacement)
                                replacements_made += 1
                                print(f"✓ Replaced '{placeholder}' with '{replacement[:30]}...'")

            elif hasattr(shape, "table"):
                table = shape.table
                for row in table.rows:
                    for cell in row.cells:
                        for paragraph in cell.text_frame.paragraphs:
                            for run in paragraph.runs:
                                for placeholder, replacement in replacements.items():
                                    if placeholder in run.text:
                                        run.text = run.text.replace(placeholder, replacement)
                                        replacements_made += 1
                                        print(f"✓ Replaced '{placeholder}' with '{replacement[:30]}...' (in table)")

        return replacements_made

    def decode_base64_image(self, base64_string):
        """
        Decode base64 image string and return image bytes
        """
        try:
            if base64_string.startswith('data:'):
                base64_string = base64_string.split(',', 1)[1]

            image_data = base64.b64decode(base64_string)

            img = Image.open(io.BytesIO(image_data))

            # Convert to RGB if necessary (for JPEG compatibility)
            if img.mode in ('RGBA', 'LA', 'P'):
                # Create white background
                background = Image.new('RGB', img.size, (255, 255, 255))
                if img.mode == 'P':
                    img = img.convert('RGBA')
                background.paste(img, mask=img.split()[-1] if img.mode in ('RGBA', 'LA') else None)
                img = background

            img_buffer = io.BytesIO()
            img.save(img_buffer, format='PNG')
            img_buffer.seek(0)

            return img_buffer.getvalue(), img.size

        except Exception as e:
            print(f"❌ Error decoding base64 image: {e}")
            return None, None

    def add_logo_from_file(self, slide, logo_path):
        """
        Add logo from file path instead of base64
        """
        try:
            if not os.path.exists(logo_path):
                print(f"❌ Logo file not found: {logo_path}")
                return False

            cm_to_inches = 0.393701

            left = Inches(29.81 * cm_to_inches)
            top = Inches(0.81 * cm_to_inches)
            width = Inches(2.87 * cm_to_inches)
            height = Inches(2.53 * cm_to_inches)

            pic = slide.shapes.add_picture(logo_path, left, top, width, height)
            print(f"✓ Added logo from {logo_path} at position ({29.81:.2f}cm, {0.81:.2f}cm)")
            return True

        except Exception as e:
            print(f"❌ Error adding logo: {e}")
            return False

    def add_logo_at_fixed_position(self, slide, logo_data):
        """
        Add logo at fixed position using binary data
        """
        try:
            cm_to_inches = 0.393701

            left = Inches(29.81 * cm_to_inches)
            top = Inches(0.81 * cm_to_inches)
            width = Inches(2.87 * cm_to_inches)
            height = Inches(2.53 * cm_to_inches)

            with tempfile.NamedTemporaryFile(suffix='.png', delete=False) as tmp_file:
                tmp_file.write(logo_data)
                tmp_logo_path = tmp_file.name

            try:
                pic = slide.shapes.add_picture(tmp_logo_path, left, top, width, height)
                print(f"✓ Added logo at position ({29.81:.2f}cm, {0.81:.2f}cm) with size {2.87:.2f}cm x {2.53:.2f}cm")
                return True

            finally:
                try:
                    os.unlink(tmp_logo_path)
                except:
                    pass

        except Exception as e:
            print(f"❌ Error adding logo: {e}")
            return False

    def add_logo_placeholder(self, slide):
        """
        Add a placeholder text box where logo should be
        """
        try:
            cm_to_inches = 0.393701

            left = Inches(29.81 * cm_to_inches)
            top = Inches(0.81 * cm_to_inches)
            width = Inches(2.87 * cm_to_inches)
            height = Inches(2.53 * cm_to_inches)

            textbox = slide.shapes.add_textbox(left, top, width, height)
            text_frame = textbox.text_frame
            text_frame.text = "Logo Here"

            paragraph = text_frame.paragraphs[0]
            paragraph.alignment = 1
            font = paragraph.font
            font.size = Inches(0.2)
            font.color.rgb = RGBColor(128, 128, 128)

            line = textbox.line
            line.color.rgb = RGBColor(200, 200, 200)
            line.width = Inches(0.01)

            print(f"✓ Added logo placeholder at position ({29.81:.2f}cm, {0.81:.2f}cm)")
            return True

        except Exception as e:
            print(f"❌ Error adding logo placeholder: {e}")
            return False

    def create_powerpoint(self, customer_name, about_client, project_name,
                          challenge_text, solution_text, impact_text,
                          logo_path=None, logo_base64=None, output_path=None):
        """
        Create PowerPoint from template by replacing placeholders

        Args:
            customer_name (str): Customer name
            about_client (str): About client text
            project_name (str): Project name
            challenge_text (str): Challenge description
            solution_text (str): Solution description
            impact_text (str): Impact description
            logo_path (str, optional): Path to logo file
            logo_base64 (str, optional): Base64 encoded logo
            output_path (str, optional): Output file path

        Returns:
            str: Path to created PowerPoint file
        """
        try:
            if not os.path.exists(self.template_path):
                raise FileNotFoundError(f"Template file not found at {self.template_path}")

            prs = Presentation(self.template_path)
            print(f"✓ Loaded template with {len(prs.slides)} slides")

            # Define replacements
            replacements = {
                '{{CUSTOMER_NAME}}': customer_name,
                '{{ABOUT_CLIENT}}': about_client,
                '{{PROJECT_NAME}}': project_name,
                '{{CHALLENGE_TEXT}}': challenge_text,
                '{{SOLUTION_TEXT}}': solution_text,
                '{{IMPACT_TEXT}}': impact_text,
            }

            # Handle logo
            has_logo = False
            logo_data = None

            if logo_base64 and logo_base64.strip():
                print("🖼️ Processing base64 logo...")
                logo_data, logo_size = self.decode_base64_image(logo_base64.strip())
                if logo_data:
                    print(f"✓ Logo decoded successfully ({len(logo_data)} bytes)")
                    has_logo = True
                else:
                    print("❌ Failed to decode base64 logo")
            elif logo_path and os.path.exists(logo_path):
                print(f"🖼️ Using logo file: {logo_path}")
                has_logo = True

            # Process slides
            for slide_idx, slide in enumerate(prs.slides):
                print(f"\n🔄 Processing slide {slide_idx + 1}...")

                # Replace text placeholders
                replacements_made = self.find_and_replace_text_in_slide(slide, replacements)
                print(f"✓ Made {replacements_made} text replacements on slide {slide_idx + 1}")

                # Add logo to first slide
                if slide_idx == 0:
                    if has_logo:
                        if logo_data:  # Base64 logo
                            self.add_logo_at_fixed_position(slide, logo_data)
                        elif logo_path:  # File logo
                            self.add_logo_from_file(slide, logo_path)
                    else:
                        self.add_logo_placeholder(slide)

            # Save the presentation
            if not output_path:
                output_path = f"reference_{customer_name.replace(' ', '_').lower()}.pptx"

            prs.save(output_path)
            print(f"\n✅ PowerPoint saved to: {output_path}")

            return output_path

        except Exception as e:
            print(f"❌ Failed to create PowerPoint: {str(e)}")
            raise


def main():
    """Command line interface"""
    parser = argparse.ArgumentParser(description='Generate PowerPoint from template')

    # Required arguments
    parser.add_argument('--customer-name', required=True, help='Customer name')
    parser.add_argument('--about-client', required=True, help='About client text')
    parser.add_argument('--project-name', required=True, help='Project name')
    parser.add_argument('--challenge-text', required=True, help='Challenge description')
    parser.add_argument('--solution-text', required=True, help='Solution description')
    parser.add_argument('--impact-text', required=True, help='Impact description')

    # Optional arguments
    parser.add_argument('--template-path', help='Path to template file')
    parser.add_argument('--logo-path', help='Path to logo image file')
    parser.add_argument('--logo-base64', help='Base64 encoded logo')
    parser.add_argument('--output-path', help='Output file path')

    args = parser.parse_args()

    generator = PowerPointGenerator(args.template_path)

    try:
        output_file = generator.create_powerpoint(
            customer_name=args.customer_name,
            about_client=args.about_client,
            project_name=args.project_name,
            challenge_text=args.challenge_text,
            solution_text=args.solution_text,
            impact_text=args.impact_text,
            logo_path=args.logo_path,
            logo_base64=args.logo_base64,
            output_path=args.output_path
        )

        print(f"\n🎉 Success! PowerPoint created: {output_file}")

    except Exception as e:
        print(f"\n❌ Error: {e}")
        return 1

    return 0


def create_presentation(customer_name, about_client, project_name,
                                challenge_text, solution_text, impact_text,
                                template_path="template.pptx", logo_path=None,
                                output_path=None):
    generator = PowerPointGenerator(template_path)

    output_file = generator.create_powerpoint(
        customer_name=customer_name,
        about_client=about_client,
        project_name=project_name,
        challenge_text=challenge_text,
        solution_text=solution_text,
        impact_text=impact_text,
        logo_path=logo_path,
        output_path=output_path
    )

    return output_file
