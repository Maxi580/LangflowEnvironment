import hashlib
import os
import time
import base64
from functools import lru_cache
import tempfile
import fitz
import zipfile
from docx import Document
import PyPDF2
import mimetypes
from pptx import Presentation
from openpyxl import load_workbook
import requests
from typing import List, Dict, Any, Optional, Tuple, Set

from .image_description_cache import ImageDescriptionCache, compute_image_hash

OLLAMA_URL = os.getenv("OLLAMA_INTERNAL_URL")
DEFAULT_EMBEDDING_MODEL = os.getenv("DEFAULT_EMBEDDING_MODEL")
DEFAULT_VISION_MODEL = os.getenv("DEFAULT_VISION_MODEL")

OLLAMA_EMBEDDINGS_ENDPOINT = os.getenv("OLLAMA_EMBEDDINGS_ENDPOINT")
OLLAMA_GENERATE_ENDPOINT = os.getenv("OLLAMA_GENERATE_ENDPOINT")
OLLAMA_TAGS_ENDPOINT = os.getenv("OLLAMA_TAGS_ENDPOINT")

MAX_IMAGE_WIDTH = os.getenv("MAX_IMAGE_WIDTH")
MAX_IMAGE_HEIGHT = os.getenv("MAX_IMAGE_HEIGHT")
JPEG_QUALITY = os.getenv("JPEG_QUALITY")
IMAGE_TIMEOUT = int(os.getenv("IMAGE_TIMEOUT"))


def get_text_embedding(text: str) -> List[float]:
    """
    Get text embedding using the configured embedding model

    Args:
        text: Text to embed

    Returns:
        List of float values representing the embedding vector

    Raises:
        ValueError: If embedding request fails or returns invalid data
    """
    url = f"{OLLAMA_URL}{OLLAMA_EMBEDDINGS_ENDPOINT}"
    model = DEFAULT_EMBEDDING_MODEL

    payload = {
        "model": model,
        "prompt": text
    }

    try:
        print(f"Requesting embedding from {url} with model: {model}")
        response = requests.post(url, json=payload, timeout=30)
        response.raise_for_status()

        result = response.json()
        print(f"Ollama API response status: {response.status_code}")

        if not isinstance(result, dict):
            raise ValueError(f"Expected dict response, got {type(result)}")

        embedding = result.get("embedding")
        if embedding is None:
            raise ValueError(f"No embedding field in response. Response keys: {list(result.keys())}")

        if not isinstance(embedding, list):
            raise ValueError(f"Expected list for embedding, got {type(embedding)}")

        if len(embedding) == 0:
            raise ValueError("Received empty embedding")

        if not all(isinstance(x, (int, float)) for x in embedding):
            raise ValueError("Embedding contains non-numeric values")

        print(f"Successfully got embedding of size {len(embedding)} for model {model}")
        return embedding

    except requests.exceptions.Timeout:
        print(f"Timeout connecting to Ollama API at {url}")
        raise ValueError(f"Timeout connecting to Ollama API (model: {model})")

    except requests.exceptions.ConnectionError:
        print(f"Connection error to Ollama API at {url}")
        raise ValueError(f"Cannot connect to Ollama API at {OLLAMA_URL} (model: {model})")

    except requests.exceptions.RequestException as e:
        print(f"Request error connecting to Ollama API: {e}")
        if hasattr(e, 'response') and e.response is not None:
            try:
                error_detail = e.response.json()
                raise ValueError(f"Ollama API error: {error_detail}")
            except:
                raise ValueError(f"Ollama API error: {e.response.text}")
        raise ValueError(f"Error connecting to Ollama API: {str(e)}")

    except ValueError:
        raise

    except Exception as e:
        print(f"Unexpected error getting embedding: {e}")
        raise ValueError(f"Unexpected error getting embedding from model {model}: {str(e)}")


def resize_image_for_vision(image_path: str) -> bytes:
    """
    Resize and optimize image for faster vision model processing

    Args:
        image_path: Path to the image file

    Returns:
        Optimized image data as bytes
    """
    try:
        with Image.open(image_path) as img:
            # Convert to RGB if necessary (handles RGBA, grayscale, etc.)
            if img.mode not in ('RGB', 'L'):
                img = img.convert('RGB')

            # Get original size
            original_size = img.size

            # Resize if image is larger than max size
            if img.size[0] > MAX_IMAGE_WIDTH or img.size[1] > MAX_IMAGE_HEIGHT:
                # Use thumbnail to maintain aspect ratio
                img.thumbnail((MAX_IMAGE_WIDTH, MAX_IMAGE_HEIGHT), Image.Resampling.LANCZOS)
                print(f"Resized image from {original_size} to {img.size}")

            # Convert to bytes with compression
            output = io.BytesIO()
            img.save(output, format='JPEG', quality=JPEG_QUALITY, optimize=True)
            return output.getvalue()

    except Exception as e:
        # Fallback: read original file if resize fails
        print(f"Warning: Could not resize image {image_path}: {e}")
        with open(image_path, "rb") as f:
            return f.read()


def get_image_description(image_path: str, prompt: str = None) -> str:
    """
    Get image description using the configured vision model

    Args:
        image_path: Path to the image file
        prompt: Custom prompt for image description (optional)

    Returns:
        String description of the image

    Raises:
        ValueError: If image processing fails
        FileNotFoundError: If image file doesn't exist
    """
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image file not found: {image_path}")

    url = f"{OLLAMA_URL}{OLLAMA_GENERATE_ENDPOINT}"
    model = DEFAULT_VISION_MODEL

    if prompt is None:
        prompt = "Describe this image in detail, including objects, people, text, colors, and setting."

    try:
        image_data = resize_image_for_vision(image_path)
        image_b64 = base64.b64encode(image_data).decode('utf-8')
    except Exception as e:
        raise ValueError(f"Failed to read image file {image_path}: {str(e)}")

    payload = {
        "model": model,
        "prompt": prompt,
        "images": [image_b64],
        "stream": False
    }

    try:
        print(f"Requesting image description from {url} with model: {model}")
        response = requests.post(url, json=payload, timeout=IMAGE_TIMEOUT)
        response.raise_for_status()

        result = response.json()
        description = result.get("response", "").strip()

        if not description:
            raise ValueError("Empty response from vision model")

        print(f"Successfully got image description of length {len(description)}")
        return description

    except requests.exceptions.Timeout:
        print(f"Timeout connecting to Ollama API at {url}")
        raise ValueError(f"Timeout connecting to Ollama API (model: {model})")

    except requests.exceptions.ConnectionError:
        print(f"Connection error to Ollama API at {url}")
        raise ValueError(f"Cannot connect to Ollama API at {OLLAMA_URL} (model: {model})")

    except requests.exceptions.RequestException as e:
        print(f"Request error connecting to Ollama API: {e}")
        if hasattr(e, 'response') and e.response is not None:
            try:
                error_detail = e.response.json()
                raise ValueError(f"Ollama API error: {error_detail}")
            except:
                raise ValueError(f"Ollama API error: {e.response.text}")
        raise ValueError(f"Error connecting to Ollama API: {str(e)}")

    except Exception as e:
        print(f"Unexpected error getting image description: {e}")
        raise ValueError(f"Unexpected error getting image description from model {model}: {str(e)}")


def get_available_models() -> Dict[str, List[str]]:
    """Get all available models from Ollama and categorize them"""
    try:
        url = f"{OLLAMA_URL}{OLLAMA_TAGS_ENDPOINT}"
        response = requests.get(url, timeout=10)
        response.raise_for_status()

        data = response.json()
        models = [model["name"] for model in data.get("models", [])]

        embedding_models = []
        vision_models = []
        chat_models = []

        for model in models:
            model_lower = model.lower()
            if "embed" in model_lower:
                embedding_models.append(model)
            elif any(vision_name in model_lower for vision_name in ["llava", "bakllava", "moondream"]):
                vision_models.append(model)
            else:
                chat_models.append(model)

        return {
            "embedding": embedding_models,
            "vision": vision_models,
            "chat": chat_models,
            "all": models
        }

    except Exception as e:
        print(f"Error getting available models: {e}")
        return {
            "embedding": [],
            "vision": [],
            "chat": [],
            "all": []
        }


def get_ollama_image_description_from_bytes(image_data: bytes) -> str:
    """
    Get image description from image bytes with caching

    Args:
        image_data: Raw image bytes

    Returns:
        String description of the image
    """
    cached_description = image_cache.get_description(image_data)
    if cached_description:
        return cached_description

    try:
        if image_data.startswith(b'\xff\xd8\xff'):
            suffix = '.jpg'
        elif image_data.startswith(b'\x89PNG\r\n\x1a\n'):
            suffix = '.png'
        elif image_data.startswith(b'GIF87a') or image_data.startswith(b'GIF89a'):
            suffix = '.gif'
        elif image_data.startswith(b'RIFF') and b'WEBP' in image_data[:12]:
            suffix = '.webp'
        else:
            suffix = '.png'

        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as temp_file:
            temp_file.write(image_data)
            temp_file_path = temp_file.name

        try:
            description = get_image_description(temp_file_path)

            if not description or description.strip() == "":
                description = "No description available for this image."
            else:
                description = description.strip()

            image_cache.store_description(image_data, description)
            print(f"Generated and cached new image description (hash: {hashlib.sha256(image_data).hexdigest()[:12]}...)")

            return description

        finally:
            try:
                os.unlink(temp_file_path)
            except OSError:
                pass

    except Exception as e:
        print(f"Error getting image description from bytes: {e}")
        error_description = "Failed to describe this image."
        image_cache.store_description(image_data, error_description)
        return error_description


def extract_images_from_pdf(file_path: str) -> List[Tuple[bytes, str]]:
    """
    Extract images from PDF using PyMuPDF

    Returns:
        List of (image_bytes, description) tuples
    """
    images = []

    try:
        doc = fitz.open(file_path)

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            image_list = page.get_images()

            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)

                    # Convert to RGB if CMYK
                    if pix.n - pix.alpha < 4:
                        img_data = pix.tobytes("png")
                        description = get_ollama_image_description_from_bytes(img_data)
                        images.append((img_data, f"[Image from page {page_num + 1}]: {description}"))
                    else:
                        # Convert CMYK to RGB
                        pix1 = fitz.Pixmap(fitz.csRGB, pix)
                        img_data = pix1.tobytes("png")
                        description = get_ollama_image_description_from_bytes(img_data)
                        images.append((img_data, f"[Image from page {page_num + 1}]: {description}"))
                        pix1 = None

                    pix = None

                except Exception as e:
                    print(f"Error extracting image {img_index} from page {page_num}: {e}")
                    continue

        doc.close()

    except Exception as e:
        print(f"Error extracting images from PDF {file_path}: {e}")

    return images


def extract_images_from_docx(file_path: str) -> List[Tuple[bytes, str]]:
    """
    Extract images from Word document (.docx)

    Returns:
        List of (image_bytes, description) tuples
    """
    images = []
    seen_hashes: Set[str] = set()

    try:
        # Word documents are also zip archives
        with zipfile.ZipFile(file_path, 'r') as zip_file:
            # Look for image files in the media directory
            for file_info in zip_file.filelist:
                if file_info.filename.startswith('word/media/'):
                    # Check if it's an image file
                    if any(file_info.filename.lower().endswith(ext) for ext in
                           ['.png', '.jpg', '.jpeg', '.gif', '.bmp', '.tiff', '.emf', '.wmf']):
                        try:
                            img_data = zip_file.read(file_info.filename)

                            if file_info.filename.lower().endswith(('.emf', '.wmf')):
                                continue

                            img_hash = compute_image_hash(img_data)
                            if img_hash in seen_hashes:
                                print(f"Skipping duplicate image: {file_info.filename}")
                                continue
                            seen_hashes.add(img_hash)

                            description = get_ollama_image_description_from_bytes(img_data)
                            images.append((img_data, f"[Word embedded image]: {description}"))

                        except Exception as e:
                            print(f"Error processing Word image {file_info.filename}: {e}")
                            continue

    except Exception as e:
        print(f"Error extracting images from Word document {file_path}: {e}")

    return images


def extract_images_from_pptx(file_path: str) -> List[Tuple[bytes, str]]:
    """
    Extract images from PowerPoint presentation

    Returns:
        List of (image_bytes, description) tuples
    """
    images = []

    try:
        presentation = Presentation(file_path)

        for slide_num, slide in enumerate(presentation.slides, 1):
            for shape in slide.shapes:
                try:
                    # Check if shape has an image
                    if hasattr(shape, "image") and shape.image:
                        img_data = shape.image.blob
                        description = get_ollama_image_description_from_bytes(img_data)
                        images.append((img_data, f"[Image from slide {slide_num}]: {description}"))

                except Exception as e:
                    print(f"Error extracting image from slide {slide_num}: {e}")
                    continue

    except Exception as e:
        print(f"Error extracting images from PowerPoint {file_path}: {e}")

    return images


def extract_images_from_xlsx(file_path: str) -> List[Tuple[bytes, str]]:
    """
    Extract images from Excel workbook

    Returns:
        List of (image_bytes, description) tuples
    """
    images = []

    try:
        # Excel files are zip archives, we can extract images directly
        with zipfile.ZipFile(file_path, 'r') as zip_file:
            # Look for image files in the media directory
            for file_info in zip_file.filelist:
                if file_info.filename.startswith('xl/media/'):
                    # Check if it's an image file
                    if any(file_info.filename.lower().endswith(ext) for ext in
                           ['.png', '.jpg', '.jpeg', '.gif', '.bmp']):
                        try:
                            img_data = zip_file.read(file_info.filename)
                            description = get_ollama_image_description_from_bytes(img_data)
                            images.append((img_data, f"[Excel embedded image]: {description}"))
                        except Exception as e:
                            print(f"Error processing Excel image {file_info.filename}: {e}")
                            continue

    except Exception as e:
        print(f"Error extracting images from Excel {file_path}: {e}")

    return images


def detect_file_type(file_path: str) -> str:
    """
    Detect file type based on extension and MIME type
    Now supports PDF, PowerPoint, Excel, and text files
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext == '.pdf':
        return 'pdf'
    elif ext == '.pptx':
        return 'pptx'
    elif ext == '.xlsx':
        return 'xlsx'
    elif ext == '.docx':
        return 'docx'
    if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
        return 'image'
    elif ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.csv']:
        return 'text'

    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type:
        if mime_type == 'application/pdf':
            return 'pdf'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return 'pptx'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return 'xlsx'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            return 'docx'
        elif mime_type.startswith('text/'):
            return 'text'

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            f.read(1024)
        return 'text'
    except UnicodeDecodeError:
        try:
            with open(file_path, "rb") as f:
                PyPDF2.PdfReader(f)
            return 'pdf'
        except:
            pass

    return 'unknown'


def extract_pdf(file_path: str, include_images: bool = True) -> str:
    """Extract text from PDF files using PyPDF2, optionally including image descriptions"""
    try:
        text = ""
        with open(file_path, "rb") as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n\n"

        if include_images:
            images = extract_images_from_pdf(file_path)
            if images:
                text += "\n\n=== EMBEDDED IMAGES ===\n\n"
                for i, (img_data, description) in enumerate(images, 1):
                    text += f"Image {i}: {description}\n\n"

        return text
    except Exception as e:
        print(f"Error extracting text from PDF {file_path}: {e}")
        raise


def extract_docx(file_path: str, include_images: bool = True) -> str:
    """
    Extract text from Word (.docx) files using python-docx, optionally including image descriptions
    """
    try:
        text_content = []
        doc = Document(file_path)

        # Extract text from paragraphs
        for paragraph in doc.paragraphs:
            if paragraph.text.strip():
                text_content.append(paragraph.text.strip())

        # Extract text from tables
        for table in doc.tables:
            table_text = ["=== TABLE ==="]
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_data.append(cell.text.strip())
                if row_data:
                    table_text.append(" | ".join(row_data))

            if len(table_text) > 1:  # More than just the header
                text_content.extend(table_text)
                text_content.append("")  # Add blank line after table

        result = "\n".join(text_content)

        # Add image descriptions if requested
        if include_images:
            images = extract_images_from_docx(file_path)
            if images:
                result += "\n\n=== EMBEDDED IMAGES ===\n\n"
                for i, (img_data, description) in enumerate(images, 1):
                    result += f"Image {i}: {description}\n\n"

        if not result.strip():
            return "No text content found in Word document."

        return result

    except Exception as e:
        print(f"Error extracting text from Word document {file_path}: {e}")
        raise ValueError(f"Failed to extract text from Word document: {str(e)}")


def extract_xlsx(file_path: str, include_images: bool = True) -> str:
    """
    Extract text from Excel (.xlsx) files using openpyxl, optionally including image descriptions
    """
    try:
        text_content = []
        workbook = load_workbook(file_path, data_only=True)

        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            sheet_text = [f"=== WORKSHEET: {sheet_name} ==="]

            max_row = worksheet.max_row
            max_col = worksheet.max_column

            first_row = 1
            for row in range(1, max_row + 1):
                if any(worksheet.cell(row, col).value is not None for col in range(1, max_col + 1)):
                    first_row = row
                    break

            last_row = max_row
            for row in range(max_row, 0, -1):
                if any(worksheet.cell(row, col).value is not None for col in range(1, max_col + 1)):
                    last_row = row
                    break

            table_data = []
            for row in range(first_row, last_row + 1):
                row_data = []
                for col in range(1, max_col + 1):
                    cell_value = worksheet.cell(row, col).value
                    if cell_value is not None:
                        cell_str = str(cell_value).strip()
                        row_data.append(cell_str)
                    else:
                        row_data.append("")

                if any(cell.strip() for cell in row_data):
                    while row_data and not row_data[-1].strip():
                        row_data.pop()

                    if row_data:
                        table_data.append(" | ".join(row_data))

            if table_data:
                sheet_text.extend(table_data)
                text_content.extend(sheet_text)
                text_content.append("")
            else:
                text_content.append(f"=== WORKSHEET: {sheet_name} ===")
                text_content.append("(Empty worksheet)")
                text_content.append("")

        workbook.close()

        result = "\n".join(text_content)

        if include_images:
            images = extract_images_from_xlsx(file_path)
            if images:
                result += "\n\n=== EMBEDDED IMAGES ===\n\n"
                for i, (img_data, description) in enumerate(images, 1):
                    result += f"Image {i}: {description}\n\n"

        if not result.strip():
            return "No data found in Excel file."

        return result

    except Exception as e:
        print(f"Error extracting text from Excel {file_path}: {e}")
        raise ValueError(f"Failed to extract text from Excel file: {str(e)}")


def extract_pptx(file_path: str, include_images: bool = True) -> str:
    """
    Extract text from PowerPoint (.pptx) files using python-pptx, optionally including image descriptions
    """
    try:
        text_content = []
        presentation = Presentation(file_path)

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = [f"=== SLIDE {slide_num} ==="]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # Handle tables specifically
                if shape.has_table:
                    table_text = []
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(" | ".join(row_text))

                    if table_text:
                        slide_text.append("TABLE:")
                        slide_text.extend(table_text)

            # Add slide content to main text
            if len(slide_text) > 1:  # More than just the slide header
                text_content.extend(slide_text)
                text_content.append("")  # Add blank line between slides

        result = "\n".join(text_content)

        # Add image descriptions if requested
        if include_images:
            images = extract_images_from_pptx(file_path)
            if images:
                result += "\n\n=== EMBEDDED IMAGES ===\n\n"
                for i, (img_data, description) in enumerate(images, 1):
                    result += f"Image {i}: {description}\n\n"

        if not result.strip():
            return "No text content found in PowerPoint presentation."

        return result

    except Exception as e:
        print(f"Error extracting text from PowerPoint {file_path}: {e}")
        raise ValueError(f"Failed to extract text from PowerPoint file: {str(e)}")


def read_file_content(file_path: str, include_images: bool = True) -> Tuple[str, str]:
    """
    Read and extract content from supported file types, optionally including image descriptions

    Args:
        file_path: Path to the file
        include_images: Whether to extract and describe embedded images

    Returns:
        (content, file_type)
    """
    file_type = detect_file_type(file_path)

    if file_type == 'text':
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        return content, 'text'
    elif file_type == 'pdf':
        content = extract_pdf(file_path, include_images)
        return content, 'pdf'
    elif file_type == 'pptx':
        content = extract_pptx(file_path, include_images)
        return content, 'pptx'
    elif file_type == 'xlsx':
        content = extract_xlsx(file_path, include_images)
        return content, 'xlsx'
    elif file_type == 'docx':
        content = extract_docx(file_path, include_images)
        return content, 'docx'
    elif file_type == 'image':
        description = get_image_description(file_path)
        return description, 'image'
    else:
        raise ValueError(f"Unsupported file type '{file_type}' for {file_path}")


image_cache = ImageDescriptionCache(max_size=1000)
