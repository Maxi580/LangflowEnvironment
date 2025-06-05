import os
import time
import base64
from functools import lru_cache
import tempfile
import fitz
import zipfile

import PyPDF2
import mimetypes
from pptx import Presentation
from openpyxl import load_workbook
import requests
from typing import List, Dict, Any, Optional, Tuple

OLLAMA_URL = os.getenv("OLLAMA_INTERNAL_URL")
DEFAULT_EMBEDDING_MODEL = os.getenv("DEFAULT_EMBEDDING_MODEL")
DEFAULT_VISION_MODEL = os.getenv("DEFAULT_VISION_MODEL")

OLLAMA_EMBEDDINGS_ENDPOINT = os.getenv("OLLAMA_EMBEDDINGS_ENDPOINT")
OLLAMA_GENERATE_ENDPOINT = os.getenv("OLLAMA_GENERATE_ENDPOINT")
OLLAMA_TAGS_ENDPOINT = os.getenv("OLLAMA_TAGS_ENDPOINT")


@lru_cache(maxsize=5)
def get_vector_size(model_name: Optional[str] = None) -> int:
    """
    Get vector size for embedding model with caching
    Only embeds sample text once per model

    Args:
        model_name: Name of the model (defaults to DEFAULT_EMBEDDING_MODEL)

    Returns:
        Vector size for the model
    """
    if model_name is None:
        model_name = DEFAULT_EMBEDDING_MODEL

    print(f"Getting vector size for model: {model_name}")

    sample_embedding = get_text_embedding("Sample text for dimension detection")
    vector_size = len(sample_embedding)

    print(f"Vector size for model '{model_name}': {vector_size}")
    return vector_size


def get_text_embedding(text: str) -> List[float]:
    """
    Get text embedding using the configured embedding model

    Args:
        text: Text to embed

    Returns:
        List of float values representing the embedding vector

    Raises:
        ValueError: If embedding request fails or returns invalid data
    """
    url = f"{OLLAMA_URL}{OLLAMA_EMBEDDINGS_ENDPOINT}"
    model = DEFAULT_EMBEDDING_MODEL

    payload = {
        "model": model,
        "prompt": text
    }

    try:
        print(f"Requesting embedding from {url} with model: {model}")
        response = requests.post(url, json=payload, timeout=30)
        response.raise_for_status()

        result = response.json()
        print(f"Ollama API response status: {response.status_code}")

        if not isinstance(result, dict):
            raise ValueError(f"Expected dict response, got {type(result)}")

        embedding = result.get("embedding")
        if embedding is None:
            raise ValueError(f"No embedding field in response. Response keys: {list(result.keys())}")

        if not isinstance(embedding, list):
            raise ValueError(f"Expected list for embedding, got {type(embedding)}")

        if len(embedding) == 0:
            raise ValueError("Received empty embedding")

        # Validate that all elements are numbers
        if not all(isinstance(x, (int, float)) for x in embedding):
            raise ValueError("Embedding contains non-numeric values")

        print(f"Successfully got embedding of size {len(embedding)} for model {model}")
        return embedding

    except requests.exceptions.Timeout:
        print(f"Timeout connecting to Ollama API at {url}")
        raise ValueError(f"Timeout connecting to Ollama API (model: {model})")

    except requests.exceptions.ConnectionError:
        print(f"Connection error to Ollama API at {url}")
        raise ValueError(f"Cannot connect to Ollama API at {OLLAMA_URL} (model: {model})")

    except requests.exceptions.RequestException as e:
        print(f"Request error connecting to Ollama API: {e}")
        if hasattr(e, 'response') and e.response is not None:
            try:
                error_detail = e.response.json()
                raise ValueError(f"Ollama API error: {error_detail}")
            except:
                raise ValueError(f"Ollama API error: {e.response.text}")
        raise ValueError(f"Error connecting to Ollama API: {str(e)}")

    except ValueError:
        raise

    except Exception as e:
        print(f"Unexpected error getting embedding: {e}")
        raise ValueError(f"Unexpected error getting embedding from model {model}: {str(e)}")


def get_image_description(image_path: str, prompt: str = None) -> str:
    """
    Get image description using the configured vision model

    Args:
        image_path: Path to the image file
        prompt: Custom prompt for image description (optional)

    Returns:
        String description of the image

    Raises:
        ValueError: If image processing fails
        FileNotFoundError: If image file doesn't exist
    """
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image file not found: {image_path}")

    url = f"{OLLAMA_URL}{OLLAMA_GENERATE_ENDPOINT}"
    model = DEFAULT_VISION_MODEL

    if prompt is None:
        prompt = "Describe this image in detail, including objects, people, text, colors, and setting."

    try:
        with open(image_path, "rb") as image_file:
            image_data = base64.b64encode(image_file.read()).decode('utf-8')
    except Exception as e:
        raise ValueError(f"Failed to read image file {image_path}: {str(e)}")

    payload = {
        "model": model,
        "prompt": prompt,
        "images": [image_data],
        "stream": False
    }

    try:
        print(f"Requesting image description from {url} with model: {model}")
        response = requests.post(url, json=payload, timeout=120)
        response.raise_for_status()

        result = response.json()
        description = result.get("response", "").strip()

        if not description:
            raise ValueError("Empty response from vision model")

        print(f"Successfully got image description of length {len(description)}")
        return description

    except requests.exceptions.Timeout:
        print(f"Timeout connecting to Ollama API at {url}")
        raise ValueError(f"Timeout connecting to Ollama API (model: {model})")

    except requests.exceptions.ConnectionError:
        print(f"Connection error to Ollama API at {url}")
        raise ValueError(f"Cannot connect to Ollama API at {OLLAMA_URL} (model: {model})")

    except requests.exceptions.RequestException as e:
        print(f"Request error connecting to Ollama API: {e}")
        if hasattr(e, 'response') and e.response is not None:
            try:
                error_detail = e.response.json()
                raise ValueError(f"Ollama API error: {error_detail}")
            except:
                raise ValueError(f"Ollama API error: {e.response.text}")
        raise ValueError(f"Error connecting to Ollama API: {str(e)}")

    except Exception as e:
        print(f"Unexpected error getting image description: {e}")
        raise ValueError(f"Unexpected error getting image description from model {model}: {str(e)}")


def test_embedding_model() -> Dict[str, Any]:
    """
    Test the configured embedding model

    Returns:
        Dictionary with test results including model info, performance metrics
    """
    try:
        model = DEFAULT_EMBEDDING_MODEL
        test_text = "This is a test sentence for embedding dimension detection."

        start_time = time.time()
        embedding = get_text_embedding(test_text)
        response_time = time.time() - start_time

        return {
            "success": True,
            "model_name": model,
            "vector_size": len(embedding),
            "response_time_seconds": round(response_time, 3),
            "sample_values": embedding[:5] if len(embedding) >= 5 else embedding,
            "test_text": test_text
        }

    except Exception as e:
        return {
            "success": False,
            "error": str(e),
            "vector_size": None,
            "model_name": DEFAULT_EMBEDDING_MODEL
        }


def test_vision_model() -> Dict[str, Any]:
    """
    Test if the configured vision model is available

    Returns:
        Dictionary with test results
    """
    model = DEFAULT_VISION_MODEL

    try:
        url = f"{OLLAMA_URL}{OLLAMA_TAGS_ENDPOINT}"
        response = requests.get(url, timeout=10)
        response.raise_for_status()

        data = response.json()
        available_models = [model["name"] for model in data.get("models", [])]

        is_available = model in available_models

        return {
            "success": is_available,
            "model_name": model,
            "available": is_available,
            "message": f"Vision model '{model}' is {'available' if is_available else 'not available'}"
        }

    except Exception as e:
        return {
            "success": False,
            "model_name": model,
            "available": False,
            "error": str(e)
        }


def get_embedding_info() -> Dict[str, Any]:
    """
    Get information about current embedding configuration

    Returns:
        Dictionary with current model configuration and capabilities
    """
    return {
        "embedding_model": DEFAULT_EMBEDDING_MODEL,
        "vision_model": DEFAULT_VISION_MODEL,
        "ollama_url": OLLAMA_URL,
        "embedding_endpoint": OLLAMA_EMBEDDINGS_ENDPOINT,
        "vision_endpoint": OLLAMA_GENERATE_ENDPOINT,
        "vector_size": get_vector_size()
    }


def clear_vector_cache():
    get_vector_size.cache_clear()
    print("Vector size cache cleared")


def get_available_models() -> Dict[str, List[str]]:
    """Get all available models from Ollama and categorize them"""
    try:
        url = f"{OLLAMA_URL}{OLLAMA_TAGS_ENDPOINT}"
        response = requests.get(url, timeout=10)
        response.raise_for_status()

        data = response.json()
        models = [model["name"] for model in data.get("models", [])]

        embedding_models = []
        vision_models = []
        chat_models = []

        for model in models:
            model_lower = model.lower()
            if "embed" in model_lower:
                embedding_models.append(model)
            elif any(vision_name in model_lower for vision_name in ["llava", "bakllava", "moondream"]):
                vision_models.append(model)
            else:
                chat_models.append(model)

        return {
            "embedding": embedding_models,
            "vision": vision_models,
            "chat": chat_models,
            "all": models
        }

    except Exception as e:
        print(f"Error getting available models: {e}")
        return {
            "embedding": [],
            "vision": [],
            "chat": [],
            "all": []
        }


def get_best_embedding_model() -> str:
    configured_model = os.getenv("DEFAULT_EMBEDDING_MODEL", "nomic-embed-text")

    available = get_available_models()
    if configured_model in available["all"]:
        return configured_model

    if available["embedding"]:
        return available["embedding"][0]

    return "nomic-embed-text"


def get_best_vision_model() -> str:
    configured_model = os.getenv("DEFAULT_VISION_MODEL", "llava:7b")

    available = get_available_models()
    if configured_model in available["all"]:
        return configured_model

    if available["vision"]:
        return available["vision"][0]

    return "llava:7b"


def get_ollama_image_description_from_bytes(image_data: bytes) -> str:
    """
    Get image description from image bytes using your existing vision model helper
    This is a wrapper around your existing get_ollama_image_description function

    Args:
        image_data: Raw image bytes

    Returns:
        String description of the image
    """
    try:
        # Determine file extension based on image data header
        # This helps Ollama better understand the image format
        if image_data.startswith(b'\xff\xd8\xff'):
            suffix = '.jpg'
        elif image_data.startswith(b'\x89PNG\r\n\x1a\n'):
            suffix = '.png'
        elif image_data.startswith(b'GIF87a') or image_data.startswith(b'GIF89a'):
            suffix = '.gif'
        elif image_data.startswith(b'RIFF') and b'WEBP' in image_data[:12]:
            suffix = '.webp'
        else:
            suffix = '.png'

        with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as temp_file:
            temp_file.write(image_data)
            temp_file_path = temp_file.name

        try:
            description = get_image_description(temp_file_path)

            if not description or description.strip() == "":
                return "No description available for this image."

            return description.strip()

        finally:
            try:
                os.unlink(temp_file_path)
            except OSError:
                pass

    except Exception as e:
        print(f"Error getting image description from bytes: {e}")
        return "Failed to describe this image."


def extract_images_from_pdf(file_path: str) -> List[Tuple[bytes, str]]:
    """
    Extract images from PDF using PyMuPDF

    Returns:
        List of (image_bytes, description) tuples
    """
    images = []

    try:
        doc = fitz.open(file_path)

        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            image_list = page.get_images()

            for img_index, img in enumerate(image_list):
                try:
                    xref = img[0]
                    pix = fitz.Pixmap(doc, xref)

                    # Convert to RGB if CMYK
                    if pix.n - pix.alpha < 4:
                        img_data = pix.tobytes("png")
                        description = get_ollama_image_description_from_bytes(img_data)
                        images.append((img_data, f"[Image from page {page_num + 1}]: {description}"))
                    else:
                        # Convert CMYK to RGB
                        pix1 = fitz.Pixmap(fitz.csRGB, pix)
                        img_data = pix1.tobytes("png")
                        description = get_ollama_image_description_from_bytes(img_data)
                        images.append((img_data, f"[Image from page {page_num + 1}]: {description}"))
                        pix1 = None

                    pix = None

                except Exception as e:
                    print(f"Error extracting image {img_index} from page {page_num}: {e}")
                    continue

        doc.close()

    except Exception as e:
        print(f"Error extracting images from PDF {file_path}: {e}")

    return images


def extract_images_from_pptx(file_path: str) -> List[Tuple[bytes, str]]:
    """
    Extract images from PowerPoint presentation

    Returns:
        List of (image_bytes, description) tuples
    """
    images = []

    try:
        presentation = Presentation(file_path)

        for slide_num, slide in enumerate(presentation.slides, 1):
            for shape in slide.shapes:
                try:
                    # Check if shape has an image
                    if hasattr(shape, "image") and shape.image:
                        img_data = shape.image.blob
                        description = get_ollama_image_description_from_bytes(img_data)
                        images.append((img_data, f"[Image from slide {slide_num}]: {description}"))

                except Exception as e:
                    print(f"Error extracting image from slide {slide_num}: {e}")
                    continue

    except Exception as e:
        print(f"Error extracting images from PowerPoint {file_path}: {e}")

    return images


def extract_images_from_xlsx(file_path: str) -> List[Tuple[bytes, str]]:
    """
    Extract images from Excel workbook

    Returns:
        List of (image_bytes, description) tuples
    """
    images = []

    try:
        # Excel files are zip archives, we can extract images directly
        with zipfile.ZipFile(file_path, 'r') as zip_file:
            # Look for image files in the media directory
            for file_info in zip_file.filelist:
                if file_info.filename.startswith('xl/media/'):
                    # Check if it's an image file
                    if any(file_info.filename.lower().endswith(ext) for ext in
                           ['.png', '.jpg', '.jpeg', '.gif', '.bmp']):
                        try:
                            img_data = zip_file.read(file_info.filename)
                            description = get_ollama_image_description_from_bytes(img_data)
                            images.append((img_data, f"[Excel embedded image]: {description}"))
                        except Exception as e:
                            print(f"Error processing Excel image {file_info.filename}: {e}")
                            continue

    except Exception as e:
        print(f"Error extracting images from Excel {file_path}: {e}")

    return images


def detect_file_type(file_path: str) -> str:
    """
    Detect file type based on extension and MIME type
    Now supports PDF, PowerPoint, Excel, and text files
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext == '.pdf':
        return 'pdf'
    elif ext == '.pptx':
        return 'pptx'
    elif ext == '.xlsx':
        return 'xlsx'
    if ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
        return 'image'
    elif ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.csv']:
        return 'text'

    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type:
        if mime_type == 'application/pdf':
            return 'pdf'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return 'pptx'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return 'xlsx'
        elif mime_type.startswith('text/'):
            return 'text'

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            f.read(1024)
        return 'text'
    except UnicodeDecodeError:
        try:
            with open(file_path, "rb") as f:
                PyPDF2.PdfReader(f)
            return 'pdf'
        except:
            pass

    return 'unknown'


def extract_pdf(file_path: str, include_images: bool = True) -> str:
    """Extract text from PDF files using PyPDF2, optionally including image descriptions"""
    try:
        text = ""
        with open(file_path, "rb") as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n\n"

        if include_images:
            images = extract_images_from_pdf(file_path)
            if images:
                text += "\n\n=== EMBEDDED IMAGES ===\n\n"
                for i, (img_data, description) in enumerate(images, 1):
                    text += f"Image {i}: {description}\n\n"

        return text
    except Exception as e:
        print(f"Error extracting text from PDF {file_path}: {e}")
        raise


def extract_xlsx(file_path: str, include_images: bool = True) -> str:
    """
    Extract text from Excel (.xlsx) files using openpyxl, optionally including image descriptions
    """
    try:
        text_content = []
        workbook = load_workbook(file_path, data_only=True)

        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            sheet_text = [f"=== WORKSHEET: {sheet_name} ==="]

            max_row = worksheet.max_row
            max_col = worksheet.max_column

            first_row = 1
            for row in range(1, max_row + 1):
                if any(worksheet.cell(row, col).value is not None for col in range(1, max_col + 1)):
                    first_row = row
                    break

            last_row = max_row
            for row in range(max_row, 0, -1):
                if any(worksheet.cell(row, col).value is not None for col in range(1, max_col + 1)):
                    last_row = row
                    break

            table_data = []
            for row in range(first_row, last_row + 1):
                row_data = []
                for col in range(1, max_col + 1):
                    cell_value = worksheet.cell(row, col).value
                    if cell_value is not None:
                        cell_str = str(cell_value).strip()
                        row_data.append(cell_str)
                    else:
                        row_data.append("")

                if any(cell.strip() for cell in row_data):
                    while row_data and not row_data[-1].strip():
                        row_data.pop()

                    if row_data:
                        table_data.append(" | ".join(row_data))

            if table_data:
                sheet_text.extend(table_data)
                text_content.extend(sheet_text)
                text_content.append("")
            else:
                text_content.append(f"=== WORKSHEET: {sheet_name} ===")
                text_content.append("(Empty worksheet)")
                text_content.append("")

        workbook.close()

        result = "\n".join(text_content)

        if include_images:
            images = extract_images_from_xlsx(file_path)
            if images:
                result += "\n\n=== EMBEDDED IMAGES ===\n\n"
                for i, (img_data, description) in enumerate(images, 1):
                    result += f"Image {i}: {description}\n\n"

        if not result.strip():
            return "No data found in Excel file."

        return result

    except Exception as e:
        print(f"Error extracting text from Excel {file_path}: {e}")
        raise ValueError(f"Failed to extract text from Excel file: {str(e)}")


def extract_pptx(file_path: str, include_images: bool = True) -> str:
    """
    Extract text from PowerPoint (.pptx) files using python-pptx, optionally including image descriptions
    """
    try:
        text_content = []
        presentation = Presentation(file_path)

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = [f"=== SLIDE {slide_num} ==="]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # Handle tables specifically
                if shape.has_table:
                    table_text = []
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(" | ".join(row_text))

                    if table_text:
                        slide_text.append("TABLE:")
                        slide_text.extend(table_text)

            # Add slide content to main text
            if len(slide_text) > 1:  # More than just the slide header
                text_content.extend(slide_text)
                text_content.append("")  # Add blank line between slides

        result = "\n".join(text_content)

        # Add image descriptions if requested
        if include_images:
            images = extract_images_from_pptx(file_path)
            if images:
                result += "\n\n=== EMBEDDED IMAGES ===\n\n"
                for i, (img_data, description) in enumerate(images, 1):
                    result += f"Image {i}: {description}\n\n"

        if not result.strip():
            return "No text content found in PowerPoint presentation."

        return result

    except Exception as e:
        print(f"Error extracting text from PowerPoint {file_path}: {e}")
        raise ValueError(f"Failed to extract text from PowerPoint file: {str(e)}")


def read_file_content(file_path: str, include_images: bool = True) -> Tuple[str, str]:
    """
    Read and extract content from supported file types, optionally including image descriptions

    Args:
        file_path: Path to the file
        include_images: Whether to extract and describe embedded images

    Returns:
        (content, file_type)
    """
    file_type = detect_file_type(file_path)

    if file_type == 'text':
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        return content, 'text'
    elif file_type == 'pdf':
        content = extract_pdf(file_path, include_images)
        return content, 'pdf'
    elif file_type == 'pptx':
        content = extract_pptx(file_path, include_images)
        return content, 'pptx'
    elif file_type == 'xlsx':
        content = extract_xlsx(file_path, include_images)
        return content, 'xlsx'
    elif file_type == 'image':
        description = get_image_description(file_path)
        return description, 'image'
    else:
        raise ValueError(f"Unsupported file type '{file_type}' for {file_path}")
