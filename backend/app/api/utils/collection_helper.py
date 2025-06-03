import os
import time
from pathlib import Path
from fastapi import Request

import requests
import uuid
from typing import List, Dict, Any, Optional, Tuple
from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams
import PyPDF2
import mimetypes
from pptx import Presentation
from openpyxl import load_workbook  # Add this import for Excel support
from ..routes.flows import get_flows

OLLAMA_URL = os.getenv("INTERNAL_OLLAMA_URL", "http://ollama:11434")
QDRANT_URL = os.getenv("INTERNAL_QDRANT_URL", "http://qdrant:6333")
DEFAULT_EMBEDDING_MODEL = os.getenv("DEFAULT_EMBEDDING_MODEL", "nomic-embed-text")
DEFAULT_CHUNK_SIZE = int(os.getenv("CHUNK_SIZE", "1000"))
DEFAULT_CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "200"))
DEFAULT_COLLECTION = os.getenv("DEFAULT_COLLECTION", "langflow_documents")


def construct_collection_name(flow_id: str):
    return flow_id


async def verify_user_flow_access(request: Request, flow_id: str) -> bool:
    try:
        print(f"Checking user access to flow: {flow_id}")
        flows = await get_flows(request, remove_example_flows=True, header_flows=False, get_all=True)
        user_flow_ids = [flow.get('id') for flow in flows if flow.get('id')]
        has_access = flow_id in user_flow_ids
        return has_access
    except Exception as e:
        print(f"Error checking flow access: {e}")
        return False


def get_ollama_embedding(text: str, model: str = DEFAULT_EMBEDDING_MODEL) -> List[float]:
    url = f"{OLLAMA_URL}/api/embeddings"
    payload = {
        "model": model,
        "prompt": text
    }

    try:
        print(f"Requesting embedding from {url} with model: {model}")
        response = requests.post(url, json=payload, timeout=30)
        response.raise_for_status()

        result = response.json()
        print(f"Ollama API response status: {response.status_code}")

        # Validate response structure
        if not isinstance(result, dict):
            raise ValueError(f"Expected dict response, got {type(result)}")

        embedding = result.get("embedding")
        if embedding is None:
            raise ValueError(f"No embedding field in response. Response keys: {list(result.keys())}")

        # Validate embedding format
        if not isinstance(embedding, list):
            raise ValueError(f"Expected list for embedding, got {type(embedding)}")

        if len(embedding) == 0:
            raise ValueError("Received empty embedding")

        # Validate that all elements are numbers
        if not all(isinstance(x, (int, float)) for x in embedding):
            raise ValueError("Embedding contains non-numeric values")

        print(f"Successfully got embedding of size {len(embedding)} for model {model}")
        return embedding

    except requests.exceptions.Timeout:
        print(f"Timeout connecting to Ollama API at {url}")
        raise ValueError(f"Timeout connecting to Ollama API (model: {model})")

    except requests.exceptions.ConnectionError:
        print(f"Connection error to Ollama API at {url}")
        raise ValueError(f"Cannot connect to Ollama API at {OLLAMA_URL} (model: {model})")

    except requests.exceptions.RequestException as e:
        print(f"Request error connecting to Ollama API: {e}")
        if hasattr(e, 'response') and e.response is not None:
            try:
                error_detail = e.response.json()
                raise ValueError(f"Ollama API error: {error_detail}")
            except:
                raise ValueError(f"Ollama API error: {e.response.text}")
        raise ValueError(f"Error connecting to Ollama API: {str(e)}")

    except ValueError:
        # Re-raise ValueError as-is
        raise

    except Exception as e:
        print(f"Unexpected error getting embedding: {e}")
        raise ValueError(f"Unexpected error getting embedding from model {model}: {str(e)}")


def test_embedding_model(model: str, base_url: str = OLLAMA_URL) -> Dict[str, Any]:
    try:
        test_text = "This is a test sentence for embedding dimension detection."

        start_time = time.time()
        embedding = get_ollama_embedding(test_text, model)
        response_time = time.time() - start_time

        return {
            "success": True,
            "model_name": model,
            "vector_size": len(embedding),
            "response_time_seconds": round(response_time, 3),
            "sample_values": embedding[:5] if len(embedding) >= 5 else embedding,
            "test_text": test_text
        }

    except Exception as e:
        return {
            "success": False,
            "model_name": model,
            "error": str(e),
            "vector_size": None
        }


def is_file_in_qdrant(file_path: str, collection_name: str, qdrant_url: str = QDRANT_URL) -> bool:
    try:
        client = QdrantClient(url=qdrant_url)
        response = client.scroll(
            collection_name=collection_name,
            scroll_filter={
                "must": [
                    {
                        "key": "metadata.file_path",
                        "match": {
                            "value": file_path
                        }
                    }
                ]
            },
            limit=1
        )
        return len(response[0]) > 0
    except Exception as e:
        print(f"Error checking if file exists in Qdrant: {e}")
        return False


def detect_file_type(file_path: str) -> str:
    """
    Detect file type based on extension and MIME type
    Now supports PDF, PowerPoint, Excel, and text files
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    # Direct extension matching
    if ext == '.pdf':
        return 'pdf'
    elif ext == '.pptx':
        return 'pptx'
    elif ext == '.xlsx':  # Add Excel support
        return 'xlsx'
    elif ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.csv']:
        return 'text'

    # MIME type fallback
    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type:
        if mime_type == 'application/pdf':
            return 'pdf'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return 'pptx'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return 'xlsx'
        elif mime_type.startswith('text/'):
            return 'text'

    # Final fallback - try to read as text
    try:
        with open(file_path, "r", encoding="utf-8") as f:
            f.read(1024)
        return 'text'
    except UnicodeDecodeError:
        try:
            with open(file_path, "rb") as f:
                PyPDF2.PdfReader(f)
            return 'pdf'
        except:
            pass

    return 'unknown'


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF files using PyPDF2"""
    try:
        text = ""
        with open(file_path, "rb") as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PDF {file_path}: {e}")
        raise


def extract_text_from_xlsx(file_path: str) -> str:
    """
    Extract text from Excel (.xlsx) files using openpyxl
    Extracts data from all worksheets, preserving table structure
    """
    try:
        text_content = []
        workbook = load_workbook(file_path, data_only=True)

        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            sheet_text = []
            sheet_text.append(f"=== WORKSHEET: {sheet_name} ===")

            max_row = worksheet.max_row
            max_col = worksheet.max_column

            first_row = 1
            for row in range(1, max_row + 1):
                if any(worksheet.cell(row, col).value is not None for col in range(1, max_col + 1)):
                    first_row = row
                    break

            # Find last non-empty row
            last_row = max_row
            for row in range(max_row, 0, -1):
                if any(worksheet.cell(row, col).value is not None for col in range(1, max_col + 1)):
                    last_row = row
                    break

            table_data = []
            for row in range(first_row, last_row + 1):
                row_data = []
                for col in range(1, max_col + 1):
                    cell_value = worksheet.cell(row, col).value
                    if cell_value is not None:
                        cell_str = str(cell_value).strip()
                        row_data.append(cell_str)
                    else:
                        row_data.append("")

                if any(cell.strip() for cell in row_data):
                    while row_data and not row_data[-1].strip():
                        row_data.pop()

                    if row_data:
                        table_data.append(" | ".join(row_data))

            if table_data:
                sheet_text.extend(table_data)
                text_content.extend(sheet_text)
                text_content.append("")
            else:
                text_content.append(f"=== WORKSHEET: {sheet_name} ===")
                text_content.append("(Empty worksheet)")
                text_content.append("")

        workbook.close()

        result = "\n".join(text_content)

        if not result.strip():
            return "No data found in Excel file."

        return result

    except Exception as e:
        print(f"Error extracting text from Excel {file_path}: {e}")
        raise ValueError(f"Failed to extract text from Excel file: {str(e)}")


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from PowerPoint (.pptx) files using python-pptx
    Extracts text from slides, text boxes, shapes, and tables
    """
    try:
        text_content = []
        presentation = Presentation(file_path)

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = []
            slide_text.append(f"=== SLIDE {slide_num} ===")

            # Extract text from all shapes in the slide
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # Handle tables specifically
                if shape.has_table:
                    table_text = []
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(" | ".join(row_text))

                    if table_text:
                        slide_text.append("TABLE:")
                        slide_text.extend(table_text)

            # Add slide content to main text
            if len(slide_text) > 1:  # More than just the slide header
                text_content.extend(slide_text)
                text_content.append("")  # Add blank line between slides

        result = "\n".join(text_content)

        if not result.strip():
            return "No text content found in PowerPoint presentation."

        return result

    except Exception as e:
        print(f"Error extracting text from PowerPoint {file_path}: {e}")
        raise ValueError(f"Failed to extract text from PowerPoint file: {str(e)}")


def read_file_content(file_path: str) -> Tuple[str, str]:
    """
    Read and extract content from supported file types
    Returns: (content, file_type)
    """
    file_type = detect_file_type(file_path)

    if file_type == 'text':
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        return content, 'text'
    elif file_type == 'pdf':
        content = extract_text_from_pdf(file_path)
        return content, 'pdf'
    elif file_type == 'pptx':
        content = extract_text_from_pptx(file_path)
        return content, 'pptx'
    elif file_type == 'xlsx':
        content = extract_text_from_xlsx(file_path)
        return content, 'xlsx'
    else:
        raise ValueError(f"Unsupported file type '{file_type}' for {file_path}")


def upload_to_qdrant(
        file_path: str,
        file_name: str,
        file_id: str,
        flow_id: str = None,
        embedding_model: str = DEFAULT_EMBEDDING_MODEL,
        chunk_size: int = DEFAULT_CHUNK_SIZE,
        chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
        qdrant_url: str = QDRANT_URL,
) -> bool:
    try:
        print(f"Processing file: {file_path}")

        collection_name = flow_id
        print(f"Using collection name: {collection_name}")

        try:
            content, file_type = read_file_content(file_path)
            print(f"File type detected: {file_type}")
            print(f"File size: {len(content)} characters")
        except ValueError as e:
            print(f"Error: {str(e)}")
            return False

        sample_embedding = get_ollama_embedding(
            "Sample text for dimension testing",
            embedding_model,
        )
        vector_size = len(sample_embedding)
        print(f"Sample embedding size: {vector_size}")

        client = QdrantClient(url=qdrant_url)

        collections = client.get_collections().collections
        collection_names = [collection.name for collection in collections]

        if collection_name not in collection_names:
            print(f"Creating new collection: {collection_name}")
            client.create_collection(
                collection_name=collection_name,
                vectors_config=VectorParams(
                    size=vector_size,
                    distance=Distance.COSINE
                )
            )
        else:
            print(f"Using existing collection: {collection_name}")

        chunks = []
        for i in range(0, len(content), chunk_size - chunk_overlap):
            chunk = content[i:i + chunk_size]
            if chunk:
                chunks.append(chunk)

        print(f"Created {len(chunks)} chunks")

        points = []
        for chunk_idx, chunk in enumerate(chunks):
            if chunk_idx % 5 == 0:
                print(f"Processing chunk {chunk_idx + 1}/{len(chunks)}")

            embedding = get_ollama_embedding(chunk, embedding_model)
            point_id = str(uuid.uuid4())

            points.append({
                "id": point_id,
                "vector": embedding,
                "payload": {
                    "page_content": chunk,
                    "metadata": {
                        "file_path": file_path,
                        "file_id": file_id,
                        "chunk_idx": chunk_idx,
                        "filename": file_name,
                        "file_type": file_type,
                        "flow_id": flow_id
                    }
                }
            })

        if points:
            client.upsert(
                collection_name=collection_name,
                points=points
            )
            print(f"Uploaded {len(points)} points to Qdrant")
            return True
        else:
            print("No points to upload")
            return False

    except Exception as e:
        print(f"Error processing file {file_path}: {e}")
        return False


def delete_file_from_qdrant(
        file_path: str,
        flow_id: str = None,
        qdrant_url: str = QDRANT_URL
) -> bool:
    try:
        collection_name = flow_id
        client = QdrantClient(url=qdrant_url)

        response = client.scroll(
            collection_name=collection_name,
            scroll_filter={
                "must": [
                    {
                        "key": "metadata.file_path",
                        "match": {
                            "value": file_path
                        }
                    }
                ]
            },
            with_payload=False,
            with_vectors=False
        )

        point_ids = [point.id for point in response[0]]

        if not point_ids:
            print(f"No points found for file: {file_path}")
            return True

        client.delete(
            collection_name=collection_name,
            points_selector=point_ids
        )

        print(f"✓ Deleted {len(point_ids)} points for file: {file_path}")
        return True

    except Exception as e:
        print(f"✗ Error deleting file from Qdrant: {e}")
        return False


def get_files_from_collection(
        collection_name: str,
        qdrant_url: str = QDRANT_URL
) -> List[Dict[str, Any]]:
    """
    Query Qdrant for all files in a specific collection

    Args:
        collection_name: Name of the collection (flow_id)
        qdrant_url: URL of Qdrant server

    Returns:
        List of file information dictionaries
    """
    try:
        client = QdrantClient(url=qdrant_url)

        collections = client.get_collections().collections
        collection_names = [collection.name for collection in collections]

        if collection_name not in collection_names:
            print(f"Collection '{collection_name}' does not exist")
            return []

        response = client.scroll(
            collection_name=collection_name,
            scroll_filter={},
            with_payload=True,
            with_vectors=False,
            limit=1000
        )

        file_info_by_path = {}
        for point in response[0]:
            if "metadata" in point.payload:
                metadata = point.payload["metadata"]
                file_path = metadata.get("file_path")

                if file_path and file_path not in file_info_by_path:
                    file_info_by_path[file_path] = {
                        "file_id": metadata.get("file_id"),
                        "file_path": file_path,
                        "file_name": metadata.get("filename"),
                        "file_type": metadata.get("file_type"),
                        "flow_id": metadata.get("flow_id")
                    }

        files = []
        for file_path, info in file_info_by_path.items():
            path_obj = Path(file_path)
            if path_obj.exists() and path_obj.is_file():
                # Add file size
                info["file_size"] = os.path.getsize(file_path)
                files.append(info)
            else:
                print(f"File not found on disk: {file_path}")

        return files

    except Exception as e:
        print(f"Error querying Qdrant for files in collection {collection_name}: {e}")
        return []