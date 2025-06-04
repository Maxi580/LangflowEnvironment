import os
import uuid
from pathlib import Path
from fastapi import Request
from typing import List, Dict, Any, Optional, Tuple

from qdrant_client import QdrantClient
from qdrant_client.models import Distance, VectorParams
import PyPDF2
import mimetypes
from pptx import Presentation
from openpyxl import load_workbook

from ..routes.flows import get_flows
from .embedding import get_text_embedding, get_image_description, get_vector_size

# Environment variables
QDRANT_URL = os.getenv("QDRANT_INTERNAL_URL")
DEFAULT_CHUNK_SIZE = int(os.getenv("CHUNK_SIZE"))
DEFAULT_CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP"))


def construct_collection_name(flow_id: str):
    """Construct collection name from flow ID"""
    return flow_id


async def verify_user_flow_access(request: Request, flow_id: str) -> bool:
    """Verify user has access to the specified flow"""
    try:
        print(f"Checking user access to flow: {flow_id}")
        flows = await get_flows(request, remove_example_flows=True, header_flows=False, get_all=True)
        user_flow_ids = [flow.get('id') for flow in flows if flow.get('id')]
        has_access = flow_id in user_flow_ids
        return has_access
    except Exception as e:
        print(f"Error checking flow access: {e}")
        return False


def is_file_in_qdrant(file_path: str, collection_name: str, qdrant_url: str = QDRANT_URL) -> bool:
    """Check if a file already exists in the Qdrant collection"""
    try:
        client = QdrantClient(url=qdrant_url)
        response = client.scroll(
            collection_name=collection_name,
            scroll_filter={
                "must": [
                    {
                        "key": "metadata.file_path",
                        "match": {
                            "value": file_path
                        }
                    }
                ]
            },
            limit=1
        )
        return len(response[0]) > 0
    except Exception as e:
        print(f"Error checking if file exists in Qdrant: {e}")
        return False


def detect_file_type(file_path: str) -> str:
    """
    Detect file type based on extension and MIME type
    Supports PDF, PowerPoint, Excel, images, and text files
    """
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    if ext == '.pdf':
        return 'pdf'
    elif ext == '.pptx':
        return 'pptx'
    elif ext == '.xlsx':
        return 'xlsx'
    elif ext in ['.jpg', '.jpeg', '.png', '.gif', '.bmp', '.tiff', '.webp']:
        return 'image'
    elif ext in ['.txt', '.md', '.py', '.js', '.html', '.css', '.json', '.xml', '.csv']:
        return 'text'

    mime_type, _ = mimetypes.guess_type(file_path)
    if mime_type:
        if mime_type == 'application/pdf':
            return 'pdf'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            return 'pptx'
        elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
            return 'xlsx'
        elif mime_type.startswith('text/'):
            return 'text'

    try:
        with open(file_path, "r", encoding="utf-8") as f:
            f.read(1024)
        return 'text'
    except UnicodeDecodeError:
        try:
            with open(file_path, "rb") as f:
                PyPDF2.PdfReader(f)
            return 'pdf'
        except:
            pass

    return 'unknown'


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text from PDF files using PyPDF2"""
    try:
        text = ""
        with open(file_path, "rb") as f:
            pdf_reader = PyPDF2.PdfReader(f)
            for page_num in range(len(pdf_reader.pages)):
                page = pdf_reader.pages[page_num]
                text += page.extract_text() + "\n\n"
        return text
    except Exception as e:
        print(f"Error extracting text from PDF {file_path}: {e}")
        raise


def extract_text_from_xlsx(file_path: str) -> str:
    """
    Extract text from Excel (.xlsx) files using openpyxl
    Extracts data from all worksheets, preserving table structure
    """
    try:
        text_content = []
        workbook = load_workbook(file_path, data_only=True)

        for sheet_name in workbook.sheetnames:
            worksheet = workbook[sheet_name]
            sheet_text = [f"=== WORKSHEET: {sheet_name} ==="]

            max_row = worksheet.max_row
            max_col = worksheet.max_column

            # Find first non-empty row
            first_row = 1
            for row in range(1, max_row + 1):
                if any(worksheet.cell(row, col).value is not None for col in range(1, max_col + 1)):
                    first_row = row
                    break

            # Find last non-empty row
            last_row = max_row
            for row in range(max_row, 0, -1):
                if any(worksheet.cell(row, col).value is not None for col in range(1, max_col + 1)):
                    last_row = row
                    break

            # Extract table data
            table_data = []
            for row in range(first_row, last_row + 1):
                row_data = []
                for col in range(1, max_col + 1):
                    cell_value = worksheet.cell(row, col).value
                    if cell_value is not None:
                        cell_str = str(cell_value).strip()
                        row_data.append(cell_str)
                    else:
                        row_data.append("")

                if any(cell.strip() for cell in row_data):
                    # Remove trailing empty cells
                    while row_data and not row_data[-1].strip():
                        row_data.pop()

                    if row_data:
                        table_data.append(" | ".join(row_data))

            if table_data:
                sheet_text.extend(table_data)
                text_content.extend(sheet_text)
                text_content.append("")
            else:
                text_content.append(f"=== WORKSHEET: {sheet_name} ===")
                text_content.append("(Empty worksheet)")
                text_content.append("")

        workbook.close()

        result = "\n".join(text_content)
        if not result.strip():
            return "No data found in Excel file."

        return result

    except Exception as e:
        print(f"Error extracting text from Excel {file_path}: {e}")
        raise ValueError(f"Failed to extract text from Excel file: {str(e)}")


def extract_text_from_pptx(file_path: str) -> str:
    """
    Extract text from PowerPoint (.pptx) files using python-pptx
    Extracts text from slides, text boxes, shapes, and tables
    """
    try:
        text_content = []
        presentation = Presentation(file_path)

        for slide_num, slide in enumerate(presentation.slides, 1):
            slide_text = [f"=== SLIDE {slide_num} ==="]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

                # Handle tables specifically
                if shape.has_table:
                    table_text = []
                    for row in shape.table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(" | ".join(row_text))

                    if table_text:
                        slide_text.append("TABLE:")
                        slide_text.extend(table_text)

            # Add slide content to main text
            if len(slide_text) > 1:  # More than just the slide header
                text_content.extend(slide_text)
                text_content.append("")  # Add blank line between slides

        result = "\n".join(text_content)
        if not result.strip():
            return "No text content found in PowerPoint presentation."

        return result

    except Exception as e:
        print(f"Error extracting text from PowerPoint {file_path}: {e}")
        raise ValueError(f"Failed to extract text from PowerPoint file: {str(e)}")


def read_file_content(file_path: str) -> Tuple[str, str]:
    """
    Read and extract content from supported file types
    Returns: (content, file_type)
    """
    file_type = detect_file_type(file_path)

    if file_type == 'text':
        with open(file_path, "r", encoding="utf-8") as f:
            content = f.read()
        return content, 'text'
    elif file_type == 'pdf':
        content = extract_text_from_pdf(file_path)
        return content, 'pdf'
    elif file_type == 'pptx':
        content = extract_text_from_pptx(file_path)
        return content, 'pptx'
    elif file_type == 'xlsx':
        content = extract_text_from_xlsx(file_path)
        return content, 'xlsx'
    elif file_type == 'image':
        description = get_image_description(file_path)
        return description, 'image'
    else:
        raise ValueError(f"Unsupported file type '{file_type}' for {file_path}")


def upload_to_qdrant(
        file_path: str,
        file_name: str,
        file_id: str,
        flow_id: str = None,
        chunk_size: int = DEFAULT_CHUNK_SIZE,
        chunk_overlap: int = DEFAULT_CHUNK_OVERLAP,
        qdrant_url: str = QDRANT_URL,
) -> bool:
    """
    Upload file content to Qdrant collection after chunking and embedding
    """
    try:
        print(f"Processing file: {file_path}")

        collection_name = flow_id
        print(f"Using collection name: {collection_name}")

        try:
            content, file_type = read_file_content(file_path)
            print(f"File type detected: {file_type}")
            print(f"File size: {len(content)} characters")
        except ValueError as e:
            print(f"Error: {str(e)}")
            return False

        vector_size = get_vector_size()
        print(f"Sample embedding size: {vector_size}")

        client = QdrantClient(url=qdrant_url)

        collections = client.get_collections().collections
        collection_names = [collection.name for collection in collections]

        if collection_name not in collection_names:
            print(f"Creating new collection: {collection_name}")
            client.create_collection(
                collection_name=collection_name,
                vectors_config=VectorParams(
                    size=vector_size,
                    distance=Distance.COSINE
                )
            )
        else:
            print(f"Using existing collection: {collection_name}")

        # Create chunks from content
        chunks = []
        for i in range(0, len(content), chunk_size - chunk_overlap):
            chunk = content[i:i + chunk_size]
            if chunk:
                chunks.append(chunk)

        print(f"Created {len(chunks)} chunks")

        # Process chunks and create points for Qdrant
        points = []
        for chunk_idx, chunk in enumerate(chunks):
            if chunk_idx % 5 == 0:
                print(f"Processing chunk {chunk_idx + 1}/{len(chunks)}")

            # Get embedding for document content
            embedding = get_text_embedding(chunk, task_type="document")
            point_id = str(uuid.uuid4())

            points.append({
                "id": point_id,
                "vector": embedding,
                "payload": {
                    "page_content": chunk,
                    "metadata": {
                        "file_path": file_path,
                        "file_id": file_id,
                        "chunk_idx": chunk_idx,
                        "filename": file_name,
                        "file_type": file_type,
                        "flow_id": flow_id
                    }
                }
            })

        # Upload to Qdrant
        if points:
            client.upsert(
                collection_name=collection_name,
                points=points
            )
            print(f"Uploaded {len(points)} points to Qdrant")
            return True
        else:
            print("No points to upload")
            return False

    except Exception as e:
        print(f"Error processing file {file_path}: {e}")
        return False


def delete_file_from_qdrant(
        file_path: str,
        flow_id: str = None,
        qdrant_url: str = QDRANT_URL
) -> bool:
    """
    Delete all points related to a specific file from Qdrant collection
    """
    try:
        collection_name = flow_id
        client = QdrantClient(url=qdrant_url)

        # Find all points for this file
        response = client.scroll(
            collection_name=collection_name,
            scroll_filter={
                "must": [
                    {
                        "key": "metadata.file_path",
                        "match": {
                            "value": file_path
                        }
                    }
                ]
            },
            with_payload=False,
            with_vectors=False
        )

        point_ids = [point.id for point in response[0]]

        if not point_ids:
            print(f"No points found for file: {file_path}")
            return True

        # Delete the points
        client.delete(
            collection_name=collection_name,
            points_selector=point_ids
        )

        print(f"✓ Deleted {len(point_ids)} points for file: {file_path}")
        return True

    except Exception as e:
        print(f"✗ Error deleting file from Qdrant: {e}")
        return False


def get_files_from_collection(
        collection_name: str,
        qdrant_url: str = QDRANT_URL
) -> List[Dict[str, Any]]:
    """
    Query Qdrant for all files in a specific collection

    Args:
        collection_name: Name of the collection (flow_id)
        qdrant_url: URL of Qdrant server

    Returns:
        List of file information dictionaries
    """
    try:
        client = QdrantClient(url=qdrant_url)

        # Check if collection exists
        collections = client.get_collections().collections
        collection_names = [collection.name for collection in collections]

        if collection_name not in collection_names:
            print(f"Collection '{collection_name}' does not exist")
            return []

        # Get all points from collection
        response = client.scroll(
            collection_name=collection_name,
            scroll_filter={},
            with_payload=True,
            with_vectors=False,
            limit=1000
        )

        # Extract unique file information
        file_info_by_path = {}
        for point in response[0]:
            if "metadata" in point.payload:
                metadata = point.payload["metadata"]
                file_path = metadata.get("file_path")

                if file_path and file_path not in file_info_by_path:
                    file_info_by_path[file_path] = {
                        "file_id": metadata.get("file_id"),
                        "file_path": file_path,
                        "file_name": metadata.get("filename"),
                        "file_type": metadata.get("file_type"),
                        "flow_id": metadata.get("flow_id")
                    }

        files = []
        for file_path, info in file_info_by_path.items():
            path_obj = Path(file_path)
            if path_obj.exists() and path_obj.is_file():
                info["file_size"] = os.path.getsize(file_path)
                files.append(info)
            else:
                print(f"File not found on disk: {file_path}")

        return files

    except Exception as e:
        print(f"Error querying Qdrant for files in collection {collection_name}: {e}")
        return []